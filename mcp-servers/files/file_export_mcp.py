# file_export_mcp.py -- Document Generation MCP Server
# Author:  Marios Adamidis (FORTHought Lab)
# Version: 1.0.0
# Generates DOCX, PPTX, XLSX, PDF, CSV, HTML, and plain text files.

import re
import os
import ast
import csv
import json
import uuid
import emoji
import math
import time
import base64
import shutil
import datetime
import tarfile
import zipfile
import py7zr
import logging
import requests
from requests import get, post
from requests.auth import HTTPBasicAuth
import threading
import markdown2
import tempfile
from pathlib import Path
from lxml import etree
from PIL import Image
from docx import Document
from docx.shared import Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.shared import qn
from docx.oxml.ns import nsdecls
from docx.oxml import parse_xml
from docx.shared import Pt as DocxPt
from bs4 import BeautifulSoup, NavigableString
from mcp.server.fastmcp import FastMCP, Context
from mcp.server.session import ServerSession
from openpyxl import Workbook, load_workbook
from openpyxl.comments import Comment
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt as PptPt
from pptx.enum.shapes import PP_PLACEHOLDER 
from pptx.parts.image import Image
from pptx.enum.text import MSO_AUTO_SIZE
from io import BytesIO
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem, Image as ReportLabImage
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.units import mm

#NonDockerImport
import asyncio
import uvicorn
from typing import Any, Union, List
from typing_extensions import TypedDict
from mcp.server.sse import SseServerTransport
from starlette.requests import Request
from starlette.applications import Starlette
from starlette.routing import Route, Mount
from starlette.responses import Response, JSONResponse, StreamingResponse

SCRIPT_VERSION = "0.8.1"

URL = os.getenv('OWUI_URL')
TOKEN = os.getenv('JWT_SECRET') ## will be deleted in 1.0.0

def _try_get_auth_header(ctx) -> str | None:
    """
    Safely extract auth header from MCP context.
    Returns None in STDIO mode (expected behavior, not an error).
    """
    try:
        req = getattr(getattr(ctx, "request_context", None), "request", None)
        if req is None:
            return None
        headers = getattr(req, "headers", None)
        if headers is None:
            return None
        return headers.get("authorization") or headers.get("Authorization")
    except Exception:
        return None


def _normalize_edits_payload(edits: dict | list | None) -> tuple[list, list, list[str]]:
    """
    Normalize various edit payload formats into canonical form.
    
    Returns:
        (edit_items, ops, warnings)
        - edit_items: list of [target_id, new_value] pairs
        - ops: list of structural operations
        - warnings: list of warning messages about format conversion
    
    Accepted formats:
        1. {"content_edits": [["pid:1", "text"]], "ops": [...]}  (canonical)
        2. {"edits": [["pid:1", "text"]], "ops": [...]}          (legacy)
        3. {"pid:1": "text", "pid:2": "text"}                    (simple dict)
        4. {"pid:1": {"text": "..."}, "pid:2": {"text": "..."}}  (model guess)
        5. [["pid:1", "text"], ...]                               (raw list)
        6. [{"target": "pid:1", "value": "text"}, ...]           (object list)
    """
    warnings = []
    ops = []
    edit_items = []
    
    if edits is None:
        return [], [], ["No edits payload provided"]
    
    # Case 1: Raw list of pairs
    if isinstance(edits, list):
        # Check if it's a list of [target, value] pairs
        if edits and isinstance(edits[0], (list, tuple)):
            return edits, [], []
        # Check if it's a list of {"target": ..., "value": ...} objects
        if edits and isinstance(edits[0], dict):
            normalized = []
            for item in edits:
                if isinstance(item, dict) and "target" in item and "value" in item:
                    normalized.append([item["target"], item["value"]])
            if normalized:
                warnings.append("Converted object list format to pairs")
                return normalized, [], warnings
        return edits, [], []
    
    if not isinstance(edits, dict):
        return [], [], [f"Invalid edits type: {type(edits).__name__}"]
    
    # Extract ops if present
    ops = edits.get("ops", []) or []
    
    # Case 2: Canonical format with content_edits
    if "content_edits" in edits:
        edit_items = edits["content_edits"]
        # Normalize if it's object format
        if isinstance(edit_items, list) and edit_items and isinstance(edit_items[0], dict):
            edit_items = [
                [item.get("target"), item.get("value")]
                for item in edit_items
                if isinstance(item, dict) and "target" in item and "value" in item
            ]
            warnings.append("Converted content_edits from object format")
        return edit_items or [], ops, warnings
    
    # Case 3: Legacy format with "edits" key containing list
    if "edits" in edits:
        edit_items = edits["edits"]
        # Handle nested {"edits": {"edits": [...]}} (double nesting)
        if isinstance(edit_items, dict) and "edits" in edit_items:
            warnings.append("Flattened double-nested edits object")
            edit_items = edit_items.get("edits", [])
        # Normalize if it's object format
        if isinstance(edit_items, list) and edit_items and isinstance(edit_items[0], dict):
            edit_items = [
                [item.get("target"), item.get("value")]
                for item in edit_items
                if isinstance(item, dict) and "target" in item and "value" in item
            ]
            warnings.append("Converted edits from object format")
        return edit_items or [], ops, warnings
    
    # Case 4: Simple dict format {"pid:1": "text"} or {"pid:1": {"text": "..."}}
    guessed_edits = []
    for key, val in edits.items():
        # Skip known structural keys
        if key in ("ops", "content_edits", "edits"):
            continue
        
        # Check if key looks like a valid target (pid:X, tid:X, nX, sid:X, A1, etc.)
        if re.match(r'^(pid:\d+|tid:\d+|n\d+|sid:\d+|[A-Z]+\d+|.*shid:\d+)$', key, re.I):
            if isinstance(val, dict) and "text" in val:
                guessed_edits.append([key, val["text"]])
            elif isinstance(val, dict) and "value" in val:
                guessed_edits.append([key, val["value"]])
            elif isinstance(val, (str, int, float, list)):
                guessed_edits.append([key, val])
    
    if guessed_edits:
        warnings.append(f"Interpreted {len(guessed_edits)} edits from simple dict format")
        return guessed_edits, ops, warnings
    
    return [], ops, ["Could not find edits in payload"]


def _env_bool(val: str | None) -> bool:
    return str(val).strip().lower() in ("1", "true", "yes", "y", "on") if val is not None else False

PERSISTENT_FILES = _env_bool(os.getenv("PERSISTENT_FILES", "false"))
FILES_DELAY = int(os.getenv("FILES_DELAY", 60)) 

EXPORT_DIR_ENV = os.getenv("FILE_EXPORT_DIR")
EXPORT_DIR = (EXPORT_DIR_ENV or r"/output").rstrip("/")
os.makedirs(EXPORT_DIR, exist_ok=True)

BASE_URL_ENV = os.getenv("FILE_EXPORT_BASE_URL")
BASE_URL = (BASE_URL_ENV or "http://localhost:9003/files").rstrip("/")

LOG_LEVEL_ENV = os.getenv("LOG_LEVEL")
LOG_FORMAT_ENV = os.getenv(
    "LOG_FORMAT", "%(asctime)s %(levelname)s %(name)s - %(message)s"
)

DOCS_TEMPLATE_PATH = os.getenv("DOCS_TEMPLATE_DIR", "/rootPath/templates")
PPTX_TEMPLATE = None
DOCX_TEMPLATE = None
XLSX_TEMPLATE = None
PPTX_TEMPLATE_PATH = None
DOCX_TEMPLATE_PATH = None
XLSX_TEMPLATE_PATH = None

if DOCS_TEMPLATE_PATH and os.path.exists(DOCS_TEMPLATE_PATH):
    logging.debug(f"Template Folder: {DOCS_TEMPLATE_PATH}")
    for root, dirs, files in os.walk(DOCS_TEMPLATE_PATH):
        for file in files:
            fpath = os.path.join(root, file)
            if file.lower().endswith(".pptx") and PPTX_TEMPLATE_PATH is None:
                PPTX_TEMPLATE_PATH = fpath
                logging.debug(f"PPTX template: {PPTX_TEMPLATE_PATH}")
            elif file.lower().endswith(".docx") and DOCX_TEMPLATE_PATH is None:
                DOCX_TEMPLATE_PATH = fpath
            elif file.lower().endswith(".xlsx") and XLSX_TEMPLATE_PATH is None:
                XLSX_TEMPLATE_PATH = fpath
    if PPTX_TEMPLATE_PATH:
        try:
            PPTX_TEMPLATE = Presentation(PPTX_TEMPLATE_PATH)
            logging.debug(f"Using PPTX template: {PPTX_TEMPLATE_PATH}")
        except Exception as e:
            logging.warning(f"PPTX template failed to load : {e}")
            PPTX_TEMPLATE = None
    else:
        logging.debug("No PPTX template found. Creation of a blank document.")
        PPTX_TEMPLATE = None

    if DOCX_TEMPLATE_PATH and os.path.exists(DOCS_TEMPLATE_PATH):
        try:
            DOCX_TEMPLATE = Document(DOCX_TEMPLATE_PATH)
            logging.debug(f"Using DOCX template: {DOCX_TEMPLATE_PATH}")
        except Exception as e:
            logging.warning(f"DOCX template failed to load : {e}")
            DOCX_TEMPLATE = None
    else:
        logging.debug("No DOCX template found. Creation of a blank document.")
        DOCX_TEMPLATE = None
    
    XLSX_TEMPLATE_PATH = os.path.join("/rootPath/templates","Default_Template.xlsx")

    if XLSX_TEMPLATE_PATH:
        try:
            XLSX_TEMPLATE = load_workbook(XLSX_TEMPLATE_PATH)
            logging.debug(f"Using XLSX template: {XLSX_TEMPLATE_PATH}")
        except Exception as e:
            logging.warning(f"Failed to load XLSX template: {e}")
            XLSX_TEMPLATE = None
    else:
        logging.debug("No XLSX template found. Creation of a blank document.")
        XLSX_TEMPLATE = None

def search_image(query):
    log.debug(f"Searching for image with query: '{query}'")
    image_source = os.getenv("IMAGE_SOURCE", "unsplash")

    if image_source == "unsplash":
        return search_unsplash(query)
    elif image_source == "local_sd":
        return search_local_sd(query)
    elif image_source == "pexels":
        return search_pexels(query)
    else:
        log.warning(f"Image source unknown : {image_source}")
        return None

def search_local_sd(query: str):
    log.debug(f"Searching for local SD image with query: '{query}'")
    SD_URL = os.getenv("LOCAL_SD_URL")
    SD_USERNAME = os.getenv("LOCAL_SD_USERNAME")
    SD_PASSWORD = os.getenv("LOCAL_SD_PASSWORD")
    DEFAULT_MODEL = os.getenv("LOCAL_SD_DEFAULT_MODEL", "sd_xl_base_1.0.safetensors")
    DEFAULT_STEPS = int(os.getenv("LOCAL_SD_STEPS", 20))
    DEFAULT_WIDTH = int(os.getenv("LOCAL_SD_WIDTH", 512))
    DEFAULT_HEIGHT = int(os.getenv("LOCAL_SD_HEIGHT", 512))
    DEFAULT_CFG_SCALE = float(os.getenv("LOCAL_SD_CFG_SCALE", 1.5))
    DEFAULT_SCHEDULER = os.getenv("LOCAL_SD_SCHEDULER", "Karras")
    DEFAULT_SAMPLE = os.getenv("LOCAL_SD_SAMPLE", "Euler a")

    if not SD_URL:
        log.warning("LOCAL_SD_URL is not defined.")
        return None

    payload = {
        "prompt": query.strip(),
        "steps": DEFAULT_STEPS,
        "width": DEFAULT_WIDTH,
        "height": DEFAULT_HEIGHT,
        "cfg_scale": DEFAULT_CFG_SCALE,
        "sampler_name": DEFAULT_SAMPLE,
        "scheduler": DEFAULT_SCHEDULER,
        "enable_hr": False,
        "hr_upscaler": "Latent",
        "seed": -1,
        "override_settings": {
            "sd_model_checkpoint": DEFAULT_MODEL
        }
    }

    try:
        url = f"{SD_URL}/sdapi/v1/txt2img"
        log.debug(f"Sending request to local SD API at {url}")
        response = requests.post(
            url,
            json=payload,
            headers={"Content-Type": "application/json"},
            auth=HTTPBasicAuth(SD_USERNAME, SD_PASSWORD),
            timeout=30
        )
        response.raise_for_status()
        data = response.json()

        images = data.get("images", [])
        if not images:
            log.warning(f"No image generated for the request : '{query}'")
            return None

        image_b64 = images[0]
        image_data = base64.b64decode(image_b64)

        folder_path = _generate_unique_folder()
        filename = f"{query.replace(' ', '_')}.png"
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)

        with open(filepath, "wb") as f:
            f.write(image_data)

        return _public_url(folder_path, filename)

    except requests.exceptions.Timeout:
        log.error(f"Timeout during generation for : '{query}'")
    except requests.exceptions.RequestException as e:
        log.error(f"Network error : {e}")
    except Exception as e:
        log.error(f"Unexpected error : {e}")

    return None

def search_unsplash(query):
    log.debug(f"Searching Unsplash for query: '{query}'")
    api_key = os.getenv("UNSPLASH_ACCESS_KEY")
    if not api_key:
        log.warning("UNSPLASH_ACCESS_KEY is not set. Cannot search for images.")
        return None
    url = "https://api.unsplash.com/search/photos"
    params = {
        "query": query,
        "per_page": 1,
        "orientation": "landscape"
    }
    headers = {"Authorization": f"Client-ID {api_key}"}
    log.debug(f"Sending request to Unsplash API")
    try:
        response = requests.get(url, params=params, headers=headers)
        log.debug(f"Unsplash API response status: {response.status_code}")
        response.raise_for_status() 
        data = response.json()
        if data.get("results"):
            image_url = data["results"][0]["urls"]["regular"]
            log.debug(f"Found image URL for '{query}': {image_url}")
            return image_url
        else:
            log.debug(f"No results found on Unsplash for query: '{query}'")
    except requests.exceptions.RequestException as e:
        log.error(f"Network error while searching image for '{query}': {e}")
    except json.JSONDecodeError as e:
        log.error(f"Error decoding JSON from Unsplash for '{query}': {e}")
    except Exception as e:
        log.error(f"Unexpected error searching image for '{query}': {e}")
    return None 

def search_pexels(query):
    log.debug(f"Searching Pexels for query: '{query}'")
    api_key = os.getenv("PEXELS_ACCESS_KEY")
    if not api_key:
        log.warning("PEXELS_ACCESS_KEY is not set. Cannot search for images.")
        return None
    url = "https://api.pexels.com/v1/search"
    params = {
        "query": query,
        "per_page": 1,
        "orientation": "landscape"
    }
    headers = {"Authorization": f"{api_key}"}
    log.debug(f"Sending request to Pexels API")
    try:
        response = requests.get(url, params=params, headers=headers)
        log.debug(f"Pexels API response status: {response.status_code}")
        response.raise_for_status() 
        data = response.json()
        if data.get("photos"):
            image_url = data["photos"][0]["src"]["large"]
            log.debug(f"Found image URL for '{query}': {image_url}")
            return image_url
        else:
            log.debug(f"No results found on Pexels for query: '{query}'")
    except requests.exceptions.RequestException as e:
        log.error(f"Network error while searching image for '{query}': {e}")
    except json.JSONDecodeError as e:
        log.error(f"Error decoding JSON from Pexels for '{query}': {e}")
    except Exception as e:
        log.error(f"Unexpected error searching image for '{query}': {e}")
    return None

def _resolve_log_level(val: str | None) -> int:
    if not val:
        return logging.INFO
    v = val.strip()
    if v.isdigit():
        try:
            return int(v)
        except ValueError:
            return logging.INFO
    return getattr(logging, v.upper(), logging.INFO)

logging.basicConfig(
    level=_resolve_log_level(LOG_LEVEL_ENV),
    format=LOG_FORMAT_ENV,
)
log = logging.getLogger("file_export_mcp")
log.setLevel(_resolve_log_level(LOG_LEVEL_ENV))
log.info("Effective LOG_LEVEL -> %s", logging.getLevelName(log.level))

class ReviewComment(TypedDict):
    index: Union[int, str]
    comment: str

mcp = FastMCP(
    name = "file_export",
    port = int(os.getenv("MCP_HTTP_PORT", "9004")),
    host = (os.getenv("MCP_HTTP_HOST", "0.0.0.0"))
)

def dynamic_font_size(content_list, max_chars=400, base_size=28, min_size=12):
    total_chars = sum(len(line) for line in content_list)
    ratio = total_chars / max_chars if max_chars > 0 else 1
    if ratio <= 1:
        return PptPt(base_size)
    else:
        new_size = int(base_size / ratio)
        return PptPt(max(min_size, new_size))

def _public_url(folder_path: str, filename: str) -> str:
    """Build a stable public URL for a generated file."""
    folder = os.path.basename(folder_path).lstrip("/")
    name = filename.lstrip("/")
    return f"{BASE_URL}/{folder}/{name}"

def _generate_unique_folder() -> str:
    folder_name = f"export_{uuid.uuid4().hex[:10]}_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}"
    folder_path = os.path.join(EXPORT_DIR, folder_name)
    os.makedirs(folder_path, exist_ok=True)
    return folder_path

def _generate_filename(folder_path: str, ext: str, filename: str = None) -> tuple[str, str]:
    if not filename:
        filename = f"export_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.{ext}"
    base, ext = os.path.splitext(filename)
    filepath = os.path.join(folder_path, filename)
    counter = 1
    while os.path.exists(filepath):
        filename = f"{base}_{counter}{ext}"
        filepath = os.path.join(folder_path, filename)
        counter += 1
    return filepath, filename

styles = getSampleStyleSheet()
styles.add(ParagraphStyle(
    name="CustomHeading1",
    parent=styles["Heading1"],
    textColor=colors.HexColor("#0A1F44"),
    fontSize=18,
    spaceAfter=16,
    spaceBefore=12,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomHeading2",
    parent=styles["Heading2"],
    textColor=colors.HexColor("#1C3F77"),
    fontSize=14,
    spaceAfter=12,
    spaceBefore=10,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomHeading3",
    parent=styles["Heading3"],
    textColor=colors.HexColor("#3A6FB0"), 
    fontSize=12,
    spaceAfter=10,
    spaceBefore=8,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomNormal",
    parent=styles["Normal"],
    fontSize=11,
    leading=14,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomListItem",
    parent=styles["Normal"],
    fontSize=11,
    leading=14,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomCode",
    parent=styles["Code"],
    fontSize=10,
    leading=12,
    fontName="Courier",
    backColor=colors.HexColor("#F5F5F5"),
    borderColor=colors.HexColor("#CCCCCC"),
    borderWidth=1,
    leftIndent=10,
    rightIndent=10,
    topPadding=5,
    bottomPadding=5
))

# ── Unicode Font Registration ──────────────────────────────────────────────
_PDF_FONT = "Helvetica"          # fallback
_PDF_FONT_BOLD = "Helvetica-Bold"
try:
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    import glob as _glob
    _djv = next(
        (p for p in [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
            "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
        ] if os.path.isfile(p)),
        None,
    )
    _djv_bold = next(
        (p for p in [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
            "/usr/share/fonts/truetype/freefont/FreeSansBold.ttf",
        ] if os.path.isfile(p)),
        None,
    )
    if _djv:
        pdfmetrics.registerFont(TTFont("UniSans", _djv))
        _PDF_FONT = "UniSans"
        if _djv_bold:
            pdfmetrics.registerFont(TTFont("UniSans-Bold", _djv_bold))
            _PDF_FONT_BOLD = "UniSans-Bold"
        else:
            _PDF_FONT_BOLD = "UniSans"
        # Rebind every custom style to the Unicode font
        for sname in ("CustomHeading1", "CustomHeading2", "CustomHeading3",
                       "CustomNormal", "CustomListItem"):
            styles[sname].fontName = _PDF_FONT
        styles["CustomCode"].fontName = "Courier"   # keep monospace for code
        log.info(f"PDF Unicode font registered: {_djv}")
except Exception as _fe:
    log.warning(f"Could not register Unicode font, using Helvetica: {_fe}")


def _sanitize_pdf_text(text: str) -> str:
    """Strip or replace characters that ReportLab cannot render.

    Converts Unicode super/subscript characters AND caret/underscore
    notation (10^9, H_2O) into ReportLab <super>/<sub> XML markup
    so PDFs render them as actual raised/lowered glyphs.
    """
    if not text:
        return ""

    import re as _re

    # ── Unicode superscript chars → collect digits for <super> tag ──
    _SUPER = {
        "\u2070": "0", "\u00b9": "1", "\u00b2": "2", "\u00b3": "3",
        "\u2074": "4", "\u2075": "5", "\u2076": "6", "\u2077": "7",
        "\u2078": "8", "\u2079": "9", "\u207b": "-", "\u207a": "+",
    }
    # ── Unicode subscript chars → collect digits for <sub> tag ──
    _SUB = {
        "\u2080": "0", "\u2081": "1", "\u2082": "2", "\u2083": "3",
        "\u2084": "4", "\u2085": "5", "\u2086": "6", "\u2087": "7",
        "\u2088": "8", "\u2089": "9", "\u208b": "-", "\u208a": "+",
    }

    # Pass 1: Convert Unicode super/sub chars → ReportLab markup
    _super_pat = _re.compile("[" + _re.escape("".join(_SUPER.keys())) + "]+")
    _sub_pat   = _re.compile("[" + _re.escape("".join(_SUB.keys())) + "]+")

    def _repl_super(m):
        inner = "".join(_SUPER.get(c, c) for c in m.group())
        return f"<super>{inner}</super>"

    def _repl_sub(m):
        inner = "".join(_SUB.get(c, c) for c in m.group())
        return f"<sub>{inner}</sub>"

    text = _super_pat.sub(_repl_super, text)
    text = _sub_pat.sub(_repl_sub, text)

    # Pass 2: Convert caret/underscore notation from models → markup
    #   10^{-9}  →  10<super>-9</super>
    #   10^-9    →  10<super>-9</super>
    #   10^9     →  10<super>9</super>
    #   H_{2}O   →  H<sub>2</sub>O
    #   H_2O     →  H<sub>2</sub>O
    # Only match when preceded by an alphanumeric (avoids mangling underscores
    # in variable_names or markdown bold **text**)
    text = _re.sub(
        r'(?<=[\w>])\^\{([^}]+)\}',       # ^{...} braced
        r'<super>\1</super>',
        text,
    )
    text = _re.sub(
        r'(?<=[\w>])\^([+-]?\d+)',         # ^N or ^-N bare
        r'<super>\1</super>',
        text,
    )
    text = _re.sub(
        r'(?<=[A-Za-z>])_\{([^}]+)\}',    # _{...} braced
        r'<sub>\1</sub>',
        text,
    )
    text = _re.sub(
        r'(?<=[A-Za-z>])_(\d+)',           # _N bare (only after letter)
        r'<sub>\1</sub>',
        text,
    )

    # ── Common problematic chars → safe replacements ──
    _MAP = {
        "\u00b5": "u",   # micro sign → u
        "\u03bc": "u",   # Greek mu → u (context: SI prefix)
        "\u00d7": "x",   # multiplication sign
        "\u2022": "-",   # bullet
        "\u25cf": "-",   # black circle
        "\u25cb": "o",   # white circle
        "\u25a0": "-",   # black square
        "\u25a1": "o",   # white square
        "\u2014": "--",  # em dash
        "\u2013": "-",   # en dash
        "\u2018": "'",   # left single quote
        "\u2019": "'",   # right single quote
        "\u201c": '"',   # left double quote
        "\u201d": '"',   # right double quote
        "\u2026": "...", # ellipsis
        "\u2264": "<=",  # less-than-or-equal
        "\u2265": ">=",  # greater-than-or-equal
        "\u00b1": "+/-", # plus-minus
        "\u2248": "~",   # approximately equal
        "\u00b0": "deg", # degree sign
        "\u212b": "A",   # angstrom sign
    }
    for k, v in _MAP.items():
        text = text.replace(k, v)

    # If no Unicode font loaded, strip remaining non-ASCII
    # but preserve our <super>/<sub> XML tags (they're pure ASCII)
    if _PDF_FONT == "Helvetica":
        text = text.encode("ascii", errors="replace").decode("ascii")
    return text


def render_text_with_emojis(text: str) -> str:
    if not text:
        return ""
    try:
        converted = emoji.emojize(text, language="alias")
        return _sanitize_pdf_text(converted)
    except Exception as e:
        log.error(f"Error in emoji conversion: {e}")
        return _sanitize_pdf_text(text)

def process_list_items(ul_or_ol_element, is_ordered=False):
    items = []
    bullet_type = '1' if is_ordered else 'bullet'
    for li in ul_or_ol_element.find_all('li', recursive=False):
        li_text_parts = []
        for content in li.contents:
            if isinstance(content, NavigableString):
                li_text_parts.append(str(content))
            elif content.name not in ['ul', 'ol']:
                 li_text_parts.append(content.get_text())
        li_text = ''.join(li_text_parts).strip()
        list_item_paragraph = None
        if li_text:
            rendered_text = render_text_with_emojis(li_text)
            list_item_paragraph = Paragraph(rendered_text, styles["CustomListItem"])
        sub_lists = li.find_all(['ul', 'ol'], recursive=False)
        sub_flowables = []
        if list_item_paragraph:
             sub_flowables.append(list_item_paragraph)
        for sub_list in sub_lists:
            is_sub_ordered = sub_list.name == 'ol'
            nested_items = process_list_items(sub_list, is_sub_ordered)
            if nested_items:
                nested_list_flowable = ListFlowable(
                    nested_items,
                    bulletType='1' if is_sub_ordered else 'bullet',
                    leftIndent=10 * mm,
                    bulletIndent=5 * mm,
                    spaceBefore=2,
                    spaceAfter=2
                )
                sub_flowables.append(nested_list_flowable)
        if sub_flowables:
            items.append(ListItem(sub_flowables))
    return items

def render_html_elements(soup):
    log.debug("Starting render_html_elements...")
    story = []
    element_count = 0
    for elem in soup.children:
        element_count += 1
        log.debug(f"Processing element #{element_count}: {type(elem)}, name={getattr(elem, 'name', 'NavigableString')}")
        if isinstance(elem, NavigableString):
            text = str(elem).strip()
            if text:
                log.debug(f"Adding Paragraph from NavigableString: {text[:50]}...")
                story.append(Paragraph(render_text_with_emojis(text), styles["CustomNormal"]))
                story.append(Spacer(1, 6))
        elif hasattr(elem, 'name'):
            tag_name = elem.name
            log.debug(f"Handling tag: <{tag_name}>")
            if tag_name == "h1":
                text = render_text_with_emojis(elem.get_text().strip())
                log.debug(f"Adding H1: {text[:50]}...")
                story.append(Paragraph(text, styles["CustomHeading1"]))
                story.append(Spacer(1, 10))
            elif tag_name == "h2":
                text = render_text_with_emojis(elem.get_text().strip())
                log.debug(f"Adding H2: {text[:50]}...")
                story.append(Paragraph(text, styles["CustomHeading2"]))
                story.append(Spacer(1, 8))
            elif tag_name == "h3":
                text = render_text_with_emojis(elem.get_text().strip())
                log.debug(f"Adding H3: {text[:50]}...")
                story.append(Paragraph(text, styles["CustomHeading3"]))
                story.append(Spacer(1, 6))
            elif tag_name == "p":
                imgs = elem.find_all("img")
                if imgs:
                    for img_tag in imgs:
                        src = img_tag.get("src")
                        alt = img_tag.get("alt", "[Image]")
                        try:
                            if src and src.startswith("http"):
                                log.debug(f"Downloading image from URL: {src}")
                                response = requests.get(src)
                                response.raise_for_status()
                                img_data = BytesIO(response.content)
                                img = ReportLabImage(img_data, width=200, height=150)
                            else:
                                log.debug(f"Loading local image: {src}")
                                img = ReportLabImage(src, width=200, height=150)
                            story.append(img)
                            story.append(Spacer(1, 10))
                        except Exception as e:
                            log.error(f"Error loading image {src}: {e}")
                            story.append(Paragraph(f"[Image: {alt}]", styles["CustomNormal"]))
                            story.append(Spacer(1, 6))
                else:
                    text = render_text_with_emojis(elem.get_text().strip())
                    if text:
                        log.debug(f"Adding Paragraph: {text[:50]}...")
                        story.append(Paragraph(text, styles["CustomNormal"]))
                        story.append(Spacer(1, 6))
            elif tag_name in ["ul", "ol"]:
                is_ordered = tag_name == "ol"
                log.debug(f"Processing list (ordered={is_ordered})...")
                items = process_list_items(elem, is_ordered)
                if items:
                    log.debug(f"Adding ListFlowable with {len(items)} items")
                    story.append(ListFlowable(items,
                        bulletType='1' if is_ordered else 'bullet',
                        leftIndent=10 * mm,
                        bulletIndent=5 * mm,
                        spaceBefore=6,
                        spaceAfter=10
                    ))
            elif tag_name == "blockquote":
                text = render_text_with_emojis(elem.get_text().strip())
                if text:
                    log.debug(f"Adding Blockquote: {text[:50]}...")
                    story.append(Paragraph(f"{text}", styles["CustomNormal"]))
                    story.append(Spacer(1, 8))
            elif tag_name in ["code", "pre"]:
                text = elem.get_text().strip()
                if text:
                    log.debug(f"Adding Code/Pre block: {text[:50]}...")
                    story.append(Paragraph(text, styles["CustomCode"]))
                    story.append(Spacer(1, 6 if tag_name == "code" else 8))
            elif tag_name == "img":
                src = elem.get("src")
                alt = elem.get("alt", "[Image]")
                log.debug(f"Found <img> tag. src='{src}', alt='{alt}'")
                if src is not None: 
                    try:
                        if src.startswith("image_query:"):

                            query = src.replace("image_query:", "").strip()
                            log.debug(f"Handling image_query: '{query}'")
                            image_url = search_image(query)
                            if image_url:
                                log.debug(f"Downloading image from Unsplash URL: {image_url}")
                                response = requests.get(image_url)
                                log.debug(f"Image download response status: {response.status_code}")
                                response.raise_for_status()
                                img_data = BytesIO(response.content)
                                img = ReportLabImage(img_data, width=200, height=150)
                                log.debug("Adding ReportLab Image object to story (Unsplash)")
                                story.append(img)
                                story.append(Spacer(1, 10))
                            else:
                                log.warning(f"No image found for query: {query}")
                                story.append(Paragraph(f"[Image non trouvee pour: {query}]", styles["CustomNormal"]))
                                story.append(Spacer(1, 6))
                        elif src.startswith("http"):
                            log.debug(f"Downloading image from direct URL: {src}")
                            response = requests.get(src)
                            log.debug(f"Image download response status: {response.status_code}")
                            response.raise_for_status()
                            img_data = BytesIO(response.content)
                            img = ReportLabImage(img_data, width=200, height=150)
                            log.debug("Adding ReportLab Image object to story (Direct URL)")
                            story.append(img)
                            story.append(Spacer(1, 10))
                        else:
                            log.debug(f"Loading local image: {src}")
                            if os.path.exists(src):
                                img = ReportLabImage(src, width=200, height=150)
                                log.debug("Adding ReportLab Image object to story (Local)")
                                story.append(img)
                                story.append(Spacer(1, 10))
                            else:
                               log.error(f"Local image file not found: {src}")
                               story.append(Paragraph(f"[Image locale non trouvee: {src}]", styles["CustomNormal"]))
                               story.append(Spacer(1, 6))
                    except requests.exceptions.RequestException as e:
                        log.error(f"Network error loading image {src}: {e}")
                        story.append(Paragraph(f"[Image (erreur reseau): {alt}]", styles["CustomNormal"]))
                        story.append(Spacer(1, 6))
                    except Exception as e:
                        log.error(f"Error processing image {src}: {e}", exc_info=True) 
                        story.append(Paragraph(f"[Image: {alt}]", styles["CustomNormal"]))
                        story.append(Spacer(1, 6))
                else:
                    log.warning("Image tag found with no 'src' attribute.")
                    story.append(Paragraph(f"[Image: {alt} (source manquante)]", styles["CustomNormal"]))
                    story.append(Spacer(1, 6))
            elif tag_name == "table":
                from reportlab.platypus import Table, TableStyle
                rows_data = []
                header_rows = 0
                thead = elem.find("thead")
                tbody = elem.find("tbody")
                # Collect header rows
                if thead:
                    for tr in thead.find_all("tr"):
                        cells = [render_text_with_emojis(td.get_text().strip()) for td in tr.find_all(["th", "td"])]
                        if cells:
                            rows_data.append(cells)
                            header_rows += 1
                # Collect body rows (or all <tr> if no thead/tbody)
                source = tbody if tbody else elem
                for tr in source.find_all("tr", recursive=(tbody is None)):
                    # Skip rows already captured from thead
                    if thead and tr.parent == thead:
                        continue
                    cells = [render_text_with_emojis(td.get_text().strip()) for td in tr.find_all(["th", "td"])]
                    if cells:
                        rows_data.append(cells)
                if rows_data:
                    # Normalize column count
                    max_cols = max(len(r) for r in rows_data)
                    rows_data = [r + [""] * (max_cols - len(r)) for r in rows_data]
                    # Wrap cells in Paragraphs for proper text handling
                    cell_style = styles["CustomNormal"]
                    header_style = ParagraphStyle("TableHeader", parent=cell_style, fontName=_PDF_FONT_BOLD, fontSize=10)
                    tbl_data = []
                    for ri, row in enumerate(rows_data):
                        st = header_style if ri < max(header_rows, 1) else cell_style
                        tbl_data.append([Paragraph(str(c), st) for c in row])
                    # Build table with auto column widths
                    t = Table(tbl_data, repeatRows=max(header_rows, 1))
                    t.setStyle(TableStyle([
                        ("BACKGROUND", (0, 0), (-1, max(header_rows, 1) - 1), colors.HexColor("#E8EEF4")),
                        ("TEXTCOLOR", (0, 0), (-1, max(header_rows, 1) - 1), colors.HexColor("#0A1F44")),
                        ("FONTSIZE", (0, 0), (-1, -1), 10),
                        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#CCCCCC")),
                        ("TOPPADDING", (0, 0), (-1, -1), 4),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                        ("LEFTPADDING", (0, 0), (-1, -1), 6),
                        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                        ("ROWBACKGROUNDS", (0, max(header_rows, 1)), (-1, -1), [colors.white, colors.HexColor("#F7F9FC")]),
                        ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ]))
                    story.append(Spacer(1, 6))
                    story.append(t)
                    story.append(Spacer(1, 10))
            elif tag_name == "br":
                log.debug("Adding Spacer for <br>")
                story.append(Spacer(1, 6))
            else:
                text = elem.get_text().strip()
                if text:
                    log.debug(f"Adding Paragraph for unknown tag <{tag_name}>: {text[:50]}...")
                    story.append(Paragraph(render_text_with_emojis(text), styles["CustomNormal"]))
                    story.append(Spacer(1, 6))
    log.debug(f"Finished render_html_elements. Story contains {len(story)} elements.")
    return story

def _cleanup_files(folder_path: str, delay_minutes: int):
    def delete_files():
        time.sleep(delay_minutes * 60)
        try:
            import shutil
            shutil.rmtree(folder_path) 
            log.debug(f"Folder {folder_path} deleted.")
        except Exception as e:
            logging.error(f"Error deleting files : {e}")
    thread = threading.Thread(target=delete_files)
    thread.start()

def _convert_markdown_to_structured(markdown_content):
    """
    Converts Markdown content into a structured format for Word
    
    Args:
        markdown_content (str): Markdown content
        
    Returns:
        list: List of objects with 'text' and 'type'
    """
    if not markdown_content or not isinstance(markdown_content, str):
        return []
    
    lines = markdown_content.split('\n')
    structured = []
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        if line.startswith('# '):
            structured.append({"text": line[2:].strip(), "type": "title"})
        elif line.startswith('## '):
            structured.append({"text": line[3:].strip(), "type": "heading"})
        elif line.startswith('### '):
            structured.append({"text": line[4:].strip(), "type": "subheading"})
        elif line.startswith('#### '):
            structured.append({"text": line[5:].strip(), "type": "subheading"})
        elif line.startswith('- '):
            structured.append({"text": line[2:].strip(), "type": "bullet"})
        elif line.startswith('* '):
            structured.append({"text": line[2:].strip(), "type": "bullet"})
        elif line.startswith('**') and line.endswith('**'):
            structured.append({"text": line[2:-2].strip(), "type": "bold"})
        else:
            structured.append({"text": line, "type": "paragraph"})
    
    return structured


def _normalize_xlsx_content(content) -> list:
    """
    Normalize various XLSX content formats to a 2D array.
    
    Accepts:
    - 2D array: [["A", "B"], [1, 2]] → pass through
    - Tab-separated strings: ["A	B", "1	2"] → split
    - List of dicts with row/col: [{row:0, col:0, text:"A"}] → convert
    - List of dicts with type: [{type:"header", text:"A"}] → extract text
    """
    if not content:
        return []
    
    if isinstance(content, list) and len(content) > 0:
        first = content[0]
        
        # Already 2D array
        if isinstance(first, (list, tuple)):
            return [[cell if cell is not None else "" for cell in row] for row in content]
        
                # Tab-separated strings
        if isinstance(first, str) and '\t' in first:
            return [row.split('\t') for row in content if isinstance(row, str)]
        
        # Dicts with row/col positions
        if isinstance(first, dict) and ('row' in first or 'col' in first):
            max_row = max((item.get('row', 0) for item in content if isinstance(item, dict)), default=0)
            max_col = max((item.get('col', 0) for item in content if isinstance(item, dict)), default=0)
            grid = [['' for _ in range(max_col + 1)] for _ in range(max_row + 1)]
            for item in content:
                if isinstance(item, dict):
                    row = item.get('row', 0)
                    col = item.get('col', 0)
                    text = item.get('text', item.get('value', ''))
                    if 0 <= row <= max_row and 0 <= col <= max_col:
                        grid[row][col] = text
            return grid
        
        # Dicts with type/text (like DOCX format)
        if isinstance(first, dict) and ('type' in first or 'text' in first):
            # Priority: extract table data blocks first — this is the most
            # common pattern from models sending structured content for XLSX.
            all_rows = []
            for item in content:
                if not isinstance(item, dict):
                    if isinstance(item, str):
                        all_rows.append([item])
                    continue
                item_type = item.get("type", "")
                if item_type == "table" and "data" in item:
                    # Table block — data is already a 2D array
                    table_data = item["data"]
                    if isinstance(table_data, list):
                        for row in table_data:
                            if isinstance(row, (list, tuple)):
                                all_rows.append([c if c is not None else "" for c in row])
                            else:
                                all_rows.append([str(row)])
                elif item_type == "title":
                    # Title block — will be used as sheet name, skip from grid
                    continue
                elif item_type == "heading":
                    # Heading — add as a merged-style label row
                    all_rows.append([item.get("text", "")])
                elif item_type == "list" and "items" in item:
                    for li in item["items"]:
                        all_rows.append([str(li)])
                else:
                    text = item.get("text", item.get("value", ""))
                    if text:
                        all_rows.append([str(text)])
            return all_rows if all_rows else []
        
        # Simple string list (single column)
        if isinstance(first, str):
            return [[str(item)] for item in content]
        
        # List of numbers
        if isinstance(first, (int, float)):
            return [[item] for item in content]
    
    # Single string with tabs and newlines
    # Single string with tabs and newlines
    if isinstance(content, str):
        if '\n' in content and '\t' in content:
            return [line.split('\t') for line in content.split('\n') if line.strip()]


def _create_excel(data: list[list[str]], filename: str, folder_path: str | None = None, title: str | None = None) -> dict:
    log.debug("Creating Excel file with optional template")
    if folder_path is None:
        folder_path = _generate_unique_folder()
    
    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "xlsx")

    if XLSX_TEMPLATE:
        try:
            log.debug("Loading XLSX template...")
            wb = load_workbook(XLSX_TEMPLATE_PATH) 
            log.debug(f"Template loaded with {len(wb.sheetnames)} sheet(s)")
        except Exception as e:
            log.warning(f"Failed to load XLSX template: {e}")
            wb = Workbook()
    else:
        log.debug("No XLSX template available, creating new workbook")
        wb = Workbook()

    ws = wb.active

    from openpyxl.utils import get_column_letter 

    
    if title:
        ws.title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).strip()[:31]
        title_cell_found = False
        for row in ws.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str) and "title" in cell.value.lower():
                    cell.value = title 
                    from openpyxl.styles import Font
                    log.debug(f"Title '{title}' replaced in the cell {get_column_letter(cell.column)}{cell.row} containing  'title'")
                    title_cell_found = True
                    break
            if title_cell_found:
                break
    
    start_row, start_col = 1, 1
    if ws.auto_filter and ws.auto_filter.ref:
        try:
            from openpyxl.utils import range_boundaries
            start_col, start_row, _, _ = range_boundaries(ws.auto_filter.ref)
        except: pass

    if not data:
        wb.save(filepath)
        return {"success": True, "filepath": filepath, "filename": filename}

    template_border = ws.cell(start_row, start_col).border
    has_borders = template_border and any([template_border.top.style, template_border.bottom.style, 
                                          template_border.left.style, template_border.right.style])
    
    for r in range(max(len(data) + 10, 50)):
        for c in range(max(len(data[0]) + 5, 20)):
            cell = ws.cell(row=start_row + r, column=start_col + c)
            
            if r < len(data) and c < len(data[0]):
                cell.value = data[r][c]
                if r == 0 and data[r][c]:  
                    from openpyxl.styles import Font
                    cell.font = Font(bold=True)
                if has_borders:  
                    from openpyxl.styles import Border
                    cell.border = Border(top=template_border.top, bottom=template_border.bottom,
                                       left=template_border.left, right=template_border.right)
            else:
                cell.value = None
                if cell.has_style:
                    from openpyxl.styles import Font, PatternFill, Border, Alignment
                    cell.font, cell.fill, cell.border, cell.alignment = Font(), PatternFill(), Border(), Alignment()

    if ws.auto_filter:
        ws.auto_filter.ref = f"{get_column_letter(start_col)}{start_row}:{get_column_letter(start_col + len(data[0]) - 1)}{start_row + len(data) - 1}"
    
    for c in range(len(data[0])):
        max_len = max(len(str(data[r][c])) for r in range(len(data)))
        ws.column_dimensions[get_column_letter(start_col + c)].width = min(max_len + 2, 150)

    wb.save(filepath)

    return {"url": _public_url(folder_path, fname), "path": filepath}
def _create_csv(data: list[list[str]], filename: str, folder_path: str | None = None) -> dict:
    log.debug("Creating CSV file")
    if folder_path is None:
        folder_path = _generate_unique_folder()

    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "csv")

    with open(filepath, "w", newline="", encoding="utf-8") as f:
        if isinstance(data, list):
            csv.writer(f).writerows(data)
        else:
            csv.writer(f).writerow([data])

    return {"url": _public_url(folder_path, fname), "path": filepath}

def _create_pdf(text: str | list[str], filename: str, folder_path: str | None = None) -> dict:    
    log.debug("Creating PDF file")
    if folder_path is None:
        folder_path = _generate_unique_folder()
    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "pdf")

    md_parts = []
    if isinstance(text, list):
        for item in text:
            if isinstance(item, str):
                md_parts.append(item)
            elif isinstance(item, dict):
                t = item.get("type", "")
                if t == "title":
                    md_parts.append(f"# {item.get('text','')}")
                elif t in ("heading", "subtitle"):
                    level = item.get("level", 2)
                    prefix = "#" * max(1, min(int(level), 4))
                    md_parts.append(f"{prefix} {item.get('text','')}")
                elif t == "paragraph":
                    md_parts.append(item.get("text",""))
                elif t == "list":
                    items = item.get("items", [])
                    md_parts.append("\n".join([f"- {x}" for x in items]))
                elif t == "table":
                    data = item.get("data", [])
                    if data and len(data) > 0:
                        header = data[0]
                        table_lines = []
                        table_lines.append(
                            "| " + " | ".join(str(c) for c in header) + " |"
                        )
                        table_lines.append(
                            "| " + " | ".join("---" for _ in header) + " |"
                        )
                        for row in data[1:]:
                            padded = list(row) + [""] * (len(header) - len(row))
                            table_lines.append(
                                "| " + " | ".join(str(c) for c in padded[:len(header)]) + " |"
                            )
                        # Join table rows with single newlines so markdown2
                        # parses them as one contiguous table block.
                        # (The outer \n\n join between md_parts is correct
                        # for separating headings/paragraphs, but breaks
                        # tables if each row is a separate md_parts entry.)
                        md_parts.append("\n".join(table_lines))
                elif t in ("image", "image_query"):
                    query = item.get("query", "")
                    if query:
                        md_parts.append(f"![Image](image_query: {query})")
                else:
                    # Unknown block type — extract text if present
                    fallback = item.get("text", "")
                    if fallback:
                        md_parts.append(str(fallback))
    else:
        md_parts = [str(text or "")]
        
    md_text = "\n\n".join(md_parts)    
   
    def replace_image_query(match):
        query = match.group(1).strip()
        image_url = search_image(query)
        return f'\n\n<img src="{image_url}" alt="Image: {query}" />\n\n' if image_url else ""

    md_text = re.sub(r'!\[[^\]]*\]\(\s*image_query:\s*([^)]+)\)', replace_image_query, md_text)
    html = markdown2.markdown(md_text, extras=['fenced-code-blocks','tables','break-on-newline','cuddled-lists'])
    soup = BeautifulSoup(html, "html.parser")
    story = render_html_elements(soup) or [Paragraph("Empty Content", styles["CustomNormal"])]

    doc = SimpleDocTemplate(filepath, topMargin=72, bottomMargin=72, leftMargin=72, rightMargin=72)
    try:
        doc.build(story)
    except Exception as e:
        log.error(f"Error building PDF {fname}: {e}", exc_info=True)
        doc2 = SimpleDocTemplate(filepath)
        doc2.build([Paragraph("Error in PDF generation", styles["CustomNormal"])])

    return {"url": _public_url(folder_path, fname), "path": filepath}

def _content_blocks_to_slides(content_blocks: list, title: str = None) -> list:
    """Convert structured content blocks into slides_data format for PPTX.

    Each heading/title block in the content starts a new slide.
    Everything between headings becomes that slide's content.
    If no headings exist, all content goes on a single slide.
    """
    if not content_blocks or not isinstance(content_blocks, list):
        return []

    slides = []
    current_slide = None

    for block in content_blocks:
        if isinstance(block, str):
            # Plain string — add as paragraph to current slide
            if current_slide is None:
                current_slide = {"title": title or "", "content": []}
            current_slide["content"].append({"type": "paragraph", "text": block})
            continue
        if not isinstance(block, dict):
            continue

        block_type = block.get("type", "")

        if block_type in ("title", "heading", "slide_title", "slide"):
            # Start a new slide
            if current_slide is not None:
                slides.append(current_slide)
            current_slide = {
                "title": block.get("text", "Untitled"),
                "content": [],
            }
            # If the block itself has sub-content, include it
            if "content" in block and isinstance(block["content"], list):
                current_slide["content"].extend(block["content"])
            if "items" in block:
                current_slide["content"].append({"type": "list", "items": block["items"]})
        else:
            # Content block — attach to current slide
            if current_slide is None:
                current_slide = {"title": title or "", "content": []}
            current_slide["content"].append(block)

    if current_slide is not None:
        slides.append(current_slide)

    return slides


def _create_presentation(slides_data: list[dict], filename: str, folder_path: str | None = None, title: str | None = None) -> dict:
    if folder_path is None:
        folder_path = _generate_unique_folder()
    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "pptx")
      
    use_template = False
    prs = None
    title_layout = None
    content_layout = None

    if PPTX_TEMPLATE:
        try:
            log.debug("Attempting to load template...")
            src = PPTX_TEMPLATE
            if hasattr(PPTX_TEMPLATE, "slides") and hasattr(PPTX_TEMPLATE, "save"):
                log.debug("Template is a Presentation object, converting to BytesIO")
                buf = BytesIO()
                PPTX_TEMPLATE.save(buf); buf.seek(0)
                src = buf

            tmp = Presentation(src)
            log.debug(f"Template loaded with {len(tmp.slides)} slides")
            if len(tmp.slides) >= 1:
                prs = tmp
                use_template = True

                title_layout = prs.slides[0].slide_layout
                content_layout = prs.slides[1].slide_layout if len(prs.slides) >= 2 else prs.slides[0].slide_layout
                log.debug("Using template layouts")

                for i in range(len(prs.slides) - 1, 0, -1):
                    rId = prs.slides._sldIdLst[i].rId 
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]
        except Exception:
            log.error(f"Error loading template: {e}")
            use_template = False
            prs = None

    if not use_template:
        log.debug("No valid template, creating new presentation with default layouts")
        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[1]

    if use_template:
        log.debug("Using template title slide")
        tslide = prs.slides[0]
        if tslide.shapes.title:
            tslide.shapes.title.text = title or ""
            for p in tslide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    title_info = next(({'size': PptPt(int(child.attrib.get('sz', 2800))/100), 'bold': child.attrib.get('b', '0') == '1'} for child in title_layout.element.iter() if 'defRPr' in child.tag.split('}')[-1] and 'sz' in child.attrib), {'size': PptPt(28), 'bold': True})

                    r.font.size = title_info['size'] 
                    r.font.bold = title_info['bold']
    else:
        log.debug("Creating new title slide")
        tslide = prs.slides.add_slide(title_layout)
        if tslide.shapes.title:
            tslide.shapes.title.text = title or ""
            for p in tslide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = PptPt(28); r.font.bold = True

    EMU_PER_IN = 914400
    slide_w_in = prs.slide_width / EMU_PER_IN
    slide_h_in = prs.slide_height / EMU_PER_IN
    log.debug(f"Slide dimensions: {slide_w_in} x {slide_h_in} inches")

    page_margin = 0.5
    gutter = 0.3

    for i, slide_data in enumerate(slides_data):
        log.debug(f"Processing slide {i+1}: {slide_data.get('title', 'Untitled')}")
        if not isinstance(slide_data, dict):
            log.warning(f"Slide data is not a dict, skipping slide {i+1}")
            continue

        slide_title = slide_data.get("title", "Untitled")
        content_list = slide_data.get("content", [])
        if not isinstance(content_list, list):
            content_list = [content_list]
        log.debug(f"Adding slide with title: '{slide_title}'")
        slide = prs.slides.add_slide(content_layout)

        if slide.shapes.title:
            slide.shapes.title.text = slide_title
            for p in slide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    title_info = next(({'size': PptPt(int(child.attrib.get('sz', 2800))/100), 'bold': child.attrib.get('b', '0') == '1'} for child in content_layout.element.iter() if 'defRPr' in child.tag.split('}')[-1] and 'sz' in child.attrib), {'size': PptPt(28), 'bold': True})

                    r.font.size = title_info['size'] 
                    r.font.bold = title_info['bold']

        content_shape = None
        try:
            for ph in slide.placeholders:
                try:
                    if ph.placeholder_format.idx == 1:
                        content_shape = ph; break
                except Exception:
                    pass
            if content_shape is None:
                for ph in slide.placeholders:
                    try:
                        if ph.placeholder_format.idx != 0:
                            content_shape = ph; break
                    except Exception:
                        pass
        except Exception:
            log.error(f"Error finding content placeholder: {e}")
            pass

        title_bottom_in = 1.0 
        if slide.shapes.title:
            try:
                title_bottom_emu = slide.shapes.title.top + slide.shapes.title.height
                title_bottom_in = max(title_bottom_emu / EMU_PER_IN, 1.0)
                title_bottom_in += 0.2
            except Exception:
                title_bottom_in = 1.2 

        if content_shape is None:

            content_shape = slide.shapes.add_textbox(Inches(page_margin), Inches(title_bottom_in), Inches(slide_w_in - 2*page_margin), Inches(slide_h_in - title_bottom_in - page_margin))
            log.debug("Creating new textbox for content")
        tf = content_shape.text_frame
        try:
            tf.clear()
        except Exception:
            log.error(f"Error clearing text frame: {e}")
            pass
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        try:
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)
        except Exception:
            pass

        content_left_in, content_top_in = page_margin, title_bottom_in
        content_width_in = slide_w_in - 2*page_margin
        content_height_in = slide_h_in - (title_bottom_in + page_margin)

        image_query = slide_data.get("image_query")
        if image_query:
            image_url = search_image(image_query)
            if image_url:
                log.debug(f"Searching for image query: '{image_query}'")
                try:
                    log.debug(f"Downloading image from URL: {image_url}")
                    response = requests.get(image_url, timeout=30)
                    response.raise_for_status()
                    image_data = response.content
                    image_stream = BytesIO(image_data)
                    pos = slide_data.get("image_position", "right")
                    size = slide_data.get("image_size", "medium")
                    if size == "small":
                        img_w_in, img_h_in = 2.0, 1.5
                    elif size == "large":
                        img_w_in, img_h_in = 4.0, 3.0
                    else:
                        img_w_in, img_h_in = 3.0, 2.0
                    log.debug(f"Image dimensions: {img_w_in} x {img_h_in} inches")

                    if pos == "left":
                        img_left_in = page_margin
                        img_top_in = title_bottom_in
                        content_left_in = img_left_in + img_w_in + gutter
                        content_top_in = title_bottom_in
                        content_width_in = max(slide_w_in - page_margin - content_left_in, 2.5)
                        content_height_in = slide_h_in - (title_bottom_in + page_margin)
                    elif pos == "right":
                        img_left_in = max(slide_w_in - page_margin - img_w_in, page_margin)
                        img_top_in = title_bottom_in
                        content_left_in = page_margin
                        content_top_in = title_bottom_in
                        content_width_in = max(img_left_in - gutter - content_left_in, 2.5)
                        content_height_in = slide_h_in - (title_bottom_in + page_margin)
                    elif pos == "top":
                        img_left_in = slide_w_in - page_margin - img_w_in
                        img_top_in = title_bottom_in
                        content_left_in = page_margin
                        content_top_in = img_top_in + img_h_in + gutter
                        content_width_in = slide_w_in - 2*page_margin
                        content_height_in = max(slide_h_in - page_margin - content_top_in, 2.0)
                    elif pos == "bottom":
                        img_left_in = slide_w_in - page_margin - img_w_in
                        img_top_in = max(slide_h_in - page_margin - img_h_in, page_margin)
                        content_left_in = page_margin
                        content_top_in = title_bottom_in
                        content_width_in = slide_w_in - 2*page_margin
                        content_height_in = max(img_top_in - gutter - content_top_in, 2.0)
                    else:
                        img_left_in = max(slide_w_in - page_margin - img_w_in, page_margin)
                        img_top_in = title_bottom_in
                        content_left_in = page_margin
                        content_top_in = title_bottom_in
                        content_width_in = max(img_left_in - gutter - content_left_in, 2.5)
                        content_height_in = slide_h_in - (title_bottom_in + page_margin)

                    slide.shapes.add_picture(image_stream, Inches(img_left_in), Inches(img_top_in), Inches(img_w_in), Inches(img_h_in))
                    log.debug(f"Image added at position: left={img_left_in}, top={img_top_in}")
                except Exception:
                    pass

        try:
            content_shape.left = Inches(content_left_in)
            content_shape.top = Inches(content_top_in)
            content_shape.width = Inches(content_width_in)
            content_shape.height = Inches(content_height_in)
        except Exception:
            pass

        approx_chars_per_in = 9.5
        approx_lines_per_in = 1.6
        safe_width = max(content_width_in, 0.1)
        safe_height = max(content_height_in, 0.1)
        est_capacity = int(safe_width * approx_chars_per_in * safe_height * approx_lines_per_in)
        font_size = dynamic_font_size(content_list, max_chars=max(est_capacity, 120), base_size=24, min_size=12)

        try:
            tf = content_shape.text_frame
        except Exception:
            try:
                tf = content_shape.text_frame
            except Exception:
                log.warning("Could not access text frame for content shape")
                continue

        if not tf.paragraphs:
            tf.add_paragraph()

        # ── Separate text-renderable items from table blocks ──
        text_items = []
        table_items = []
        for item in content_list:
            if isinstance(item, dict) and item.get("type") == "table":
                table_items.append(item)
            else:
                text_items.append(item)

        # ── Render text items into the text frame ──
        _p_idx = 0
        for item in text_items:
            p = tf.paragraphs[0] if _p_idx == 0 else tf.add_paragraph()

            if isinstance(item, dict):
                block_type = item.get("type", "")

                if block_type == "paragraph":
                    run = p.add_run()
                    run.text = item.get("text", "")
                    run.font.size = font_size
                    _p_idx += 1

                elif block_type == "list":
                    for li_idx, li_text in enumerate(item.get("items", [])):
                        lp = p if li_idx == 0 else tf.add_paragraph()
                        run = lp.add_run()
                        run.text = f"\u2022  {li_text}"
                        run.font.size = font_size
                        lp.space_after = PptPt(4)
                        _p_idx += 1
                    continue  # space_after already set per bullet

                elif block_type in ("heading", "subtitle", "subheading"):
                    run = p.add_run()
                    run.text = item.get("text", "")
                    _heading_size = PptPt(font_size.pt + 4) if hasattr(font_size, "pt") else PptPt(20)
                    run.font.size = _heading_size
                    run.font.bold = True
                    _p_idx += 1

                else:
                    # Unknown dict — try to extract text
                    run = p.add_run()
                    run.text = item.get("text", str(item))
                    run.font.size = font_size
                    _p_idx += 1
            else:
                # Plain string or other
                run = p.add_run()
                run.text = str(item) if item is not None else ""
                run.font.size = font_size
                _p_idx += 1

            p.space_after = PptPt(6)

        # ── Render table blocks as actual PPTX table shapes ──
        for table_item in table_items:
            table_data = table_item.get("data", [])
            if not table_data:
                continue
            n_rows = len(table_data)
            n_cols = max((len(r) for r in table_data if isinstance(r, (list, tuple))), default=0)
            if n_rows == 0 or n_cols == 0:
                continue

            slide_h_emu = int(slide_h_in * EMU_PER_IN)
            margin_emu = int(page_margin * EMU_PER_IN)
            row_height_emu = int(0.4 * EMU_PER_IN)
            tbl_height = min(n_rows * row_height_emu, int((slide_h_in - 2) * EMU_PER_IN))
            gap_emu = int(0.15 * EMU_PER_IN)

            # If text items exist, shrink the text box to make room for the table
            if text_items:
                available_emu = slide_h_emu - margin_emu  # bottom margin
                needed_for_table = tbl_height + gap_emu
                max_text_height = available_emu - int(title_bottom_in * EMU_PER_IN) - needed_for_table
                if max_text_height > int(0.5 * EMU_PER_IN):
                    try:
                        content_shape.height = max_text_height
                    except Exception:
                        pass
                tbl_top_emu = content_shape.top + content_shape.height + gap_emu
            else:
                # No text — place table right below the title
                tbl_top_emu = int(title_bottom_in * EMU_PER_IN) + gap_emu

            # Clamp: never go below slide bottom minus margin
            max_top = slide_h_emu - margin_emu - tbl_height
            tbl_top_emu = min(tbl_top_emu, max(max_top, margin_emu))

            tbl_left = int(page_margin * EMU_PER_IN)
            tbl_width = int((slide_w_in - 2 * page_margin) * EMU_PER_IN)

            tbl_shape = slide.shapes.add_table(
                n_rows, n_cols,
                tbl_left, tbl_top_emu, tbl_width, tbl_height
            )
            tbl = tbl_shape.table
            tbl_font_size = PptPt(max(font_size.pt - 2, 10)) if hasattr(font_size, "pt") else PptPt(11)

            for ri, row in enumerate(table_data):
                if not isinstance(row, (list, tuple)):
                    continue
                for ci, val in enumerate(row):
                    if ci >= n_cols:
                        break
                    cell = tbl.cell(ri, ci)
                    cell.text = str(val) if val is not None else ""
                    for para in cell.text_frame.paragraphs:
                        para.font.size = tbl_font_size
                        if ri == 0:
                            para.font.bold = True

    prs.save(filepath)
    return {"url": _public_url(folder_path, fname), "path": filepath}

def _create_word(content: list[dict] | str, filename: str, folder_path: str | None = None, title: str | None = None) -> dict:
    log.debug("Creating Word document")

    if isinstance(content, str):
        content = _convert_markdown_to_structured(content)
    elif not isinstance(content, list):
        content = []

    if folder_path is None:
        folder_path = _generate_unique_folder()
    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "docx")

    use_template = False
    doc = None

    if DOCX_TEMPLATE:
        try:
            src = DOCX_TEMPLATE
            if hasattr(DOCX_TEMPLATE, "paragraphs") and hasattr(DOCX_TEMPLATE, "save"):
                buf = BytesIO()
                DOCX_TEMPLATE.save(buf)
                buf.seek(0)
                src = buf

            doc = Document(src)
            use_template = True
            log.debug("Using DOCX template")

            for element in doc.element.body:
                if element.tag.endswith('}p') or element.tag.endswith('}tbl'):
                    doc.element.body.remove(element)

        except Exception as e:
            log.warning(f"Failed to load DOCX template: {e}")
            use_template = False
            doc = None

    if not use_template:
        doc = Document()
        log.debug("Creating new Word document without template")

    if title:
        title_paragraph = doc.add_paragraph(title)
        try:
            title_paragraph.style = doc.styles['Title']
        except KeyError:
            try:
                title_paragraph.style = doc.styles['Heading 1']
            except KeyError:
                run = title_paragraph.runs[0] if title_paragraph.runs else title_paragraph.add_run()
                run.font.size = DocxPt(20)
                run.font.bold = True
        title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        log.debug("Document title added")

    for item in content or []:
        if isinstance(item, str):
            doc.add_paragraph(item)
        elif isinstance(item, dict):
            if item.get("type") == "image_query":
                new_item = {
                    "type": "image",
                    "query": item.get("query")
                }
                image_query = new_item.get("query")
                if image_query:
                    log.debug(f"Image search for the query : {image_query}")
                    image_url = search_image(image_query)
                    if image_url:
                        response = requests.get(image_url)
                        image_data = BytesIO(response.content)
                        doc.add_picture(image_data, width=Inches(6))
                        log.debug("Image successfully added")
                    else:
                        log.warning(f"Image search for : '{image_query}'")
            elif "type" in item:
                item_type = item.get("type")
                if item_type == "title":
                    paragraph = doc.add_paragraph(item.get("text", ""))
                    try:
                        paragraph.style = doc.styles['Heading 1']
                    except KeyError:
                        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                        run.font.size = DocxPt(18)
                        run.font.bold = True
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    log.debug("Title added")
                elif item_type == "subtitle":
                    paragraph = doc.add_paragraph(item.get("text", ""))
                    try:
                        paragraph.style = doc.styles['Heading 2']
                    except KeyError:
                        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                        run.font.size = DocxPt(16)
                        run.font.bold = True
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    log.debug("Subtitle added")
                elif item_type == "paragraph":
                    doc.add_paragraph(item.get("text", ""))
                    log.debug("Paragraph added")
                elif item_type == "list":
                    items = item.get("items", [])
                    for i, item_text in enumerate(items):
                        paragraph = doc.add_paragraph(item_text)
                        try:
                            paragraph.style = doc.styles['List Bullet']
                        except KeyError:
                            paragraph.style = doc.styles['Normal']
                    log.debug("List added")
                elif item_type == "image":
                    image_query = item.get("query")
                    if image_query:
                        log.debug(f"Image search for the query : {image_query}")
                        image_url = search_image(image_query)
                        if image_url:
                            response = requests.get(image_url)
                            image_data = BytesIO(response.content)
                            doc.add_picture(image_data, width=Inches(6))
                            log.debug("Image successfully added")
                        else:
                            log.warning(f"Image search for : '{image_query}'")
                elif item_type == "table":
                    data = item.get("data", [])
                    if data:
                        template_table_style = None
                        if use_template and DOCX_TEMPLATE:
                            try:
                                for table in DOCX_TEMPLATE.tables:
                                    if table.style:
                                        template_table_style = table.style
                                        break
                            except Exception:
                                pass
                        
                        table = doc.add_table(rows=len(data), cols=len(data[0]) if data else 0)
                        
                        if template_table_style:
                            try:
                                table.style = template_table_style
                                log.debug(f"Applied template table style: {template_table_style.name}")
                            except Exception as e:
                                log.debug(f"Could not apply template table style: {e}")
                        else:
                            try:
                                for style_name in ['Table Grid', 'Light Grid Accent 1', 'Medium Grid 1 Accent 1', 'Light List Accent 1']:
                                    try:
                                        table.style = doc.styles[style_name]
                                        log.debug(f"Applied built-in table style: {style_name}")
                                        break
                                    except KeyError:
                                        continue
                            except Exception as e:
                                log.debug(f"Could not apply any table style: {e}")
                        
                        for i, row in enumerate(data):
                            for j, cell in enumerate(row):
                                cell_obj = table.cell(i, j)
                                cell_obj.text = str(cell)
                                
                                if i == 0:
                                    for paragraph in cell_obj.paragraphs:
                                        for run in paragraph.runs:
                                            run.font.bold = True
                                        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        
                        if not template_table_style:
                            try:
                                tbl = table._tbl
                                tblPr = tbl.tblPr
                                tblBorders = parse_xml(r'<w:tblBorders {}><w:top w:val="single" w:sz="4" w:space="0" w:color="000000"/><w:left w:val="single" w:sz="4" w:space="0" w:color="000000"/><w:bottom w:val="single" w:sz="4" w:space="0" w:color="000000"/><w:right w:val="single" w:sz="4" w:space="0" w:color="000000"/><w:insideH w:val="single" w:sz="4" w:space="0" w:color="000000"/><w:insideV w:val="single" w:sz="4" w:space="0" w:color="000000"/></w:tblBorders>'.format(nsdecls('w')))
                                tblPr.append(tblBorders)
                            except Exception as e:
                                log.debug(f"Could not add table borders: {e}")
                        
                        log.debug("Table added with improved styling")
            elif "text" in item:
                doc.add_paragraph(item["text"])
                log.debug("Paragraph added")
    
    doc.save(filepath)
    return {"url": _public_url(folder_path, fname), "path": filepath}

def _create_raw_file(content: str, filename: str | None, folder_path: str | None = None) -> dict:
    log.debug("Creating raw file")
    if folder_path is None:
        folder_path = _generate_unique_folder()

    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "txt")

    if fname.lower().endswith(".xml") and isinstance(content, str) and not content.strip().startswith("<?xml"):
        content = f'<?xml version="1.0" encoding="UTF-8"?>\n{content}'

    with open(filepath, "w", encoding="utf-8") as f:
        f.write(content or "")

    return {"url": _public_url(folder_path, fname), "path": filepath}

def upload_file(file_path: str, filename: str, file_type: str, token: str = None) -> dict:
    """
    Move the edited file to the public fileserver directory instead of 
    uploading to Open WebUI (avoids 500 errors).
    """
    try:
        # 1. Generate a unique folder in /data/files (mapped to EXPORT_DIR)
        folder_path = _generate_unique_folder()
        
        # 2. Ensure filename has extension
        if not filename.endswith(f".{file_type}"):
            filename = f"{filename}.{file_type}"
            
        final_path = os.path.join(folder_path, filename)
        
        # 3. Move the file from temp to public storage
        import shutil
        shutil.move(file_path, final_path)
        
        # 4. Generate the public URL
        # BASE_URL from FILE_EXPORT_BASE_URL env var
        folder_name = os.path.basename(folder_path)
        url = f"{BASE_URL}/{folder_name}/{filename}"
        
        log.info(f"File published locally to: {url}")
        
        # Return the structure Open WebUI expects, but with our link
        return {
            "file_path_download": f"[Download Edited File]({url})",
            "url": url,
            "id": "local_edit", # Dummy ID to satisfy some clients
            "meta": {"url": url}
        }
    except Exception as e:
        log.error(f"Failed to publish edited file: {e}")
        return {"error": {"message": str(e)}}

def download_file(file_id: str, token: str = None) -> BytesIO | dict:
    """
    Download a file from Open WebUI (disk fast-path → HTTP fallback).
    
    Returns:
        BytesIO on success
        dict with "error" key on failure
    """
    base_dir = "/openwebui_uploads"
    
    # === DISK FAST-PATH (no auth needed) ===
    if os.path.exists(base_dir):
        try:
            # A) Flat layout: /openwebui_uploads/<file_id>_<original_filename>
            import glob
            flat_matches = glob.glob(os.path.join(base_dir, f"{file_id}_*"))
            flat_matches = [p for p in flat_matches if os.path.isfile(p)]
            if flat_matches:
                file_path = sorted(flat_matches)[0]
                logging.info(f"[disk-first] Reading: {file_path}")
                with open(file_path, "rb") as f:
                    return BytesIO(f.read())

            # B) Legacy dir layout: /openwebui_uploads/<file_id>/<somefile>
            id_dir = os.path.join(base_dir, file_id)
            if os.path.isdir(id_dir):
                files = [f for f in os.listdir(id_dir) if os.path.isfile(os.path.join(id_dir, f))]
                if files:
                    file_path = os.path.join(id_dir, sorted(files)[0])
                    logging.info(f"[disk-first] Reading (legacy): {file_path}")
                    with open(file_path, "rb") as f:
                        return BytesIO(f.read())
                        
        except Exception as e:
            logging.warning(f"[disk-first] Read failed: {e}")

    # === HTTP FALLBACK (only if token provided) ===
    if not token:
        # No token and disk failed = cannot proceed
        return {"error": {"message": f"File {file_id} not found on disk and no auth token for HTTP fallback"}}
    
    logging.info(f"[http-fallback] Downloading {file_id}")
    url = f"{URL}/api/v1/files/{file_id}/content"
    
    # Normalize token format
    auth_header = token if token.startswith("Bearer ") else f"Bearer {token}"
    headers = {
        "Authorization": auth_header,
        "Accept": "application/json"
    }
    
    try:
        response = get(url, headers=headers, timeout=30)
        if response.status_code != 200:
            return {"error": {"message": f"HTTP {response.status_code} downloading file {file_id}"}}
        return BytesIO(response.content)
    except Exception as e:
        return {"error": {"message": f"HTTP download failed: {e}"}}

def _extract_paragraph_style_info(para):
    """Extrait les informations de style détaillées d'un paragraphe"""
    if not para.runs:
        return {}
    
    first_run = para.runs[0]
    return {
        "font_name": first_run.font.name,
        "font_size": first_run.font.size,
        "bold": first_run.font.bold,
        "italic": first_run.font.italic,
        "underline": first_run.font.underline,
        "color": first_run.font.color.rgb if first_run.font.color else None
    }

def _extract_cell_style_info(cell):
    """Extrait les informations de style d'une cellule"""
    return {
        "style": cell.style.name if hasattr(cell, 'style') else None,
        "text_alignment": cell.paragraphs[0].alignment if cell.paragraphs else None
    }

def _apply_text_to_paragraph(para, new_text):
    """
    Apply new text to a paragraph while preserving formatting.
    """
    original_style = para.style
    original_alignment = para.alignment
    
    original_run_format = None
    if para.runs:
        first_run = para.runs[0]
        original_run_format = {
            "font_name": first_run.font.name,
            "font_size": first_run.font.size,
            "bold": first_run.font.bold,
            "italic": first_run.font.italic,
            "underline": first_run.font.underline,
            "color": first_run.font.color.rgb if first_run.font.color and first_run.font.color.rgb else None
        }
    
    for _ in range(len(para.runs)):
        para._element.remove(para.runs[0]._element)
    
    if isinstance(new_text, list):
        for i, text_item in enumerate(new_text):
            if i > 0:
                para.add_run("\n")
            run = para.add_run(str(text_item))
            if original_run_format:
                _apply_run_formatting(run, original_run_format)
    else:
        run = para.add_run(str(new_text))
        if original_run_format:
            _apply_run_formatting(run, original_run_format)
    
    if original_style:
        try:
            para.style = original_style
        except Exception:
            pass
    if original_alignment is not None:
        try:
            para.alignment = original_alignment
        except Exception:
            pass


def _apply_run_formatting(run, format_dict):
    """
    Apply formatting from a dict to a run.
    """
    try:
        if format_dict.get("font_name"):
            run.font.name = format_dict["font_name"]
    except Exception:
        pass
    
    try:
        if format_dict.get("font_size"):
            run.font.size = format_dict["font_size"]
    except Exception:
        pass
    
    try:
        if format_dict.get("bold") is not None:
            run.font.bold = format_dict["bold"]
    except Exception:
        pass
    
    try:
        if format_dict.get("italic") is not None:
            run.font.italic = format_dict["italic"]
    except Exception:
        pass
    
    try:
        if format_dict.get("underline") is not None:
            run.font.underline = format_dict["underline"]
    except Exception:
        pass
    
    try:
        if format_dict.get("color"):
            from docx.shared import RGBColor
            run.font.color.rgb = format_dict["color"]
    except Exception:
        pass

@mcp.tool(
    name="full_context_document",
    title="Return the structure of a document (docx, xlsx, pptx)",
    description="Return the structure, content, and metadata of a document based on its type (docx, xlsx, pptx). Unified output format with index, type, style, and text."
)

async def full_context_document(
    file_id: str,
    file_name: str,
    ctx: Context[ServerSession, None]
) -> dict:
    """
    Return the structure of a document (docx, xlsx, pptx) based on its file extension.
    The function detects the file type and processes it accordingly.
    Returns:
        dict: A JSON object with the structure of the document.
    """
    bearer_token = _try_get_auth_header(ctx)
    if bearer_token:
        logging.debug("Using auth header from request context")
        user_token = bearer_token
    else:
        # STDIO mode - no header available, disk-first will handle it
        logging.debug("No auth header (STDIO mode) - relying on disk-first download")
        user_token = None
        logging.error(f"Error retrieving authorization header use admin fallback")
    try:
        user_file = download_file(file_id=file_id,token=user_token)

        if isinstance(user_file, dict) and "error" in user_file:
            return json.dumps(user_file, indent=4, ensure_ascii=False)

        file_extension = os.path.splitext(file_name)[1].lower()
        file_type = file_extension.lstrip('.')

        structure = {
            "file_name": file_name,
            "file_id": file_id,
            "type": file_type,
            "slide_id_order": [],
            "body": [],
        }
        index_counter = 1

        if file_type == "docx":
            doc = Document(user_file)
            
            para_id_counter = 1
            
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                
                style = para.style.name
                style_info = _extract_paragraph_style_info(para)
                element_type = "heading" if style.startswith("Heading") else "paragraph"
                
                para_xml_id = para_id_counter
                
                structure["body"].append({
                    "index": para_id_counter,
                    "para_xml_id": para_xml_id,
                    "id_key": f"pid:{para_xml_id}",
                    "type": element_type,
                    "style": style,
                    "style_info": style_info,
                    "text": text
                })
                para_id_counter += 1
            
            for table_idx, table in enumerate(doc.tables):
                table_xml_id = id(table._element)
                table_info = {
                    "index": para_id_counter,
                    "table_xml_id": table_xml_id,
                    "id_key": f"tid:{table_xml_id}",
                    "type": "table",
                    "style": "Table",
                    "table_id": table_idx,
                    "rows": len(table.rows),
                    "columns": len(table.rows[0].cells) if table.rows else 0,
                    "cells": []
                }
                
                for row_idx, row in enumerate(table.rows):
                    row_data = []
                    for col_idx, cell in enumerate(row.cells):
                        cell_xml_id = id(cell._element)
                        cell_text = cell.text.strip()
                        cell_data = {
                            "row": row_idx,
                            "column": col_idx,
                            "cell_xml_id": cell_xml_id,
                            "id_key": f"tid:{table_xml_id}/cid:{cell_xml_id}",
                            "text": cell_text,
                            "style": cell.style.name if hasattr(cell, 'style') else None
                        }
                        row_data.append(cell_data)
                    table_info["cells"].append(row_data)
                
                structure["body"].append(table_info)
                para_id_counter += 1

        elif file_type == "xlsx":
            wb = load_workbook(user_file, read_only=True, data_only=True)

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                for row_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                    for col_idx, cell in enumerate(row, start=1):
                        if cell is None or str(cell).strip() == "":
                            continue
                        col_letter = sheet.cell(row=row_idx, column=col_idx).column_letter
                        cell_ref = f"{col_letter}{row_idx}"
                        structure["body"].append({
                            "index": cell_ref,
                            "type": "cell",
                            "text": str(cell)
                        })
                        index_counter += 1

        elif file_type == "pptx":
            prs = Presentation(user_file)
            structure["slide_id_order"] = [int(s.slide_id) for s in prs.slides]
            for slide_idx, slide in enumerate(prs.slides):
                title_shape = slide.shapes.title if hasattr(slide.shapes, "title") else None
                title_text = title_shape.text.strip() if (title_shape and getattr(title_shape, "text", "").strip()) else ""

                slide_obj = {
                    "index": slide_idx,
                    "slide_id": int(slide.slide_id),
                    "id_key": f"sid:{int(slide.slide_id)}",
                    "title": title_text,
                    "shapes": []
                }
                

                for shape_idx, shape in enumerate(slide.shapes):
                    key = f"s{slide_idx}/sh{shape_idx}"
                    if hasattr(shape, "image"):
                        shape_id_val = getattr(shape, "shape_id", None) or shape._element.cNvPr.id  
                        slide_obj["shapes"].append({
                            "shape_idx": shape_idx,
                            "shape_id": shape_id_val,
                            "idx_key": key, 
                            "id_key": f"sid:{int(slide.slide_id)}/shid:{int(shape_id_val)}",
                            "kind": "image"
                        })
                        continue

                    if hasattr(shape, "text_frame") and shape.text_frame:
                        kind = "title" if (title_shape is not None and shape is title_shape) else "textbox"

                        paragraphs = []
                        for p in shape.text_frame.paragraphs:
                            text = "".join(run.text for run in p.runs) if p.runs else p.text
                            text = (text or "").strip()
                            if text != "":
                                paragraphs.append(text)

                        shape_id_val = getattr(shape, "shape_id", None) or shape._element.cNvPr.id 
                        slide_obj["shapes"].append({
                            "shape_idx": shape_idx,
                            "shape_id": shape_id_val,
                            "idx_key": key, 
                            "id_key": f"sid:{int(slide.slide_id)}/shid:{int(shape_id_val)}",
                            "kind": kind,
                            "paragraphs": paragraphs
                        })
                        continue
                    if getattr(shape, "has_table", False):
                        tbl = shape.table
                        rows = []
                        for r in tbl.rows:
                            row_cells = []
                            for c in r.cells:
                                # collect full text
                                if hasattr(c, "text_frame") and c.text_frame:
                                    paras = []
                                    for p in c.text_frame.paragraphs:
                                        t = "".join(run.text for run in p.runs) if p.runs else p.text
                                        t = (t or "").strip()
                                        if t:
                                            paras.append(t)
                                    cell_text = "\n".join(paras)
                                else:
                                    cell_text = (getattr(c, "text", "") or "").strip()
                                row_cells.append(cell_text)
                            rows.append(row_cells)

                        shape_id_val = getattr(shape, "shape_id", None) or shape._element.cNvPr.id
                        slide_obj["shapes"].append({
                            "shape_idx": shape_idx,
                            "shape_id": shape_id_val,
                            "idx_key": key,
                            "id_key": f"sid:{int(slide.slide_id)}/shid:{int(shape_id_val)}",
                            "kind": "table",
                            "rows": rows  # list of lists: each inner list = one row's cell texts
                        })
                        continue

                structure["body"].append(slide_obj)

        else:
            return json.dumps({
                "error": {"message": f"Unsupported file type: {file_type}. Only docx, xlsx, and pptx are supported."}
            }, indent=4, ensure_ascii=False)

        return json.dumps(structure, indent=4, ensure_ascii=False)

    except Exception as e:
        return json.dumps({"error": {"message": str(e)}}, indent=4, ensure_ascii=False)

def add_auto_sized_review_comment(cell, text, author="AI Reviewer"):
    """
    Adds a note to an Excel cell, adjusting the width and height
    so that all the text is visible.
    """
    if not text:
        return

    avg_char_width = 7
    px_per_line = 15
    base_width = 200
    max_width = 500
    min_height = 40

    width = min(max_width, base_width + len(text) * 2)
    chars_per_line = max(1, width // avg_char_width)
    lines = 0
    for paragraph in text.split('\n'):
        lines += math.ceil(len(paragraph) / chars_per_line)
    height = max(min_height, lines * px_per_line)

    comment = Comment(text, author)
    comment.width = width
    comment.height = height
    cell.comment = comment

def _snapshot_runs(p):
    """Return a list of {'text': str, 'font': {...}} for each run in a paragraph."""
    runs = []
    for r in p.runs:
        f = r.font
        font_spec = {
            "name": f.name,
            "size": f.size,
            "bold": f.bold,
            "italic": f.italic,
            "underline": f.underline,
            "color_rgb": getattr(getattr(f.color, "rgb", None), "rgb", None) or getattr(f.color, "rgb", None)
        }
        runs.append({"text": r.text or "", "font": font_spec})
    return runs

def _apply_font(run, font_spec):
    """Apply font specifications to a run."""
    if not font_spec:
        return
    f = run.font
    try:
        if font_spec.get("name") is not None:
            f.name = font_spec["name"]
        if font_spec.get("size") is not None:
            f.size = font_spec["size"]
        if font_spec.get("bold") is not None:
            f.bold = font_spec["bold"]
        if font_spec.get("italic") is not None:
            f.italic = font_spec["italic"]
        if font_spec.get("underline") is not None:
            f.underline = font_spec["underline"]
        rgb = font_spec.get("color_rgb")
        if rgb is not None:
            try:
                f.color.rgb = rgb
            except Exception:
                pass
    except Exception:
        pass

def _set_text_with_runs(shape, new_content):
    """
    Set the text of a shape while preserving the original run-level formatting.
    """

    if not (hasattr(shape, "text_frame") and shape.text_frame):
        return
    tf = shape.text_frame

    if isinstance(new_content, list):
        lines = [str(item) for item in new_content]
    else:
        lines = [str(new_content or "")]

    original_para_styles = []
    original_para_runs = []     

    for p in tf.paragraphs:
        level = int(getattr(p, "level", 0) or 0)
        alignment = getattr(p, "alignment", None)
        original_para_styles.append({"level": level, "alignment": alignment})
        original_para_runs.append(_snapshot_runs(p))

    tf.clear()

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if (i == 0 and tf.paragraphs) else tf.add_paragraph()

        if original_para_styles:
            style = original_para_styles[i] if i < len(original_para_styles) else original_para_styles[-1]
            p.level = style.get("level", 0)
            if style.get("alignment") is not None:
                p.alignment = style["alignment"]

        runs_spec = (
            original_para_runs[i] if i < len(original_para_runs)
            else (original_para_runs[-1] if original_para_runs else [])
        )

        if not runs_spec:
            r = p.add_run()
            r.text = line
            continue

        n = len(runs_spec)
        total = len(line)

        if total == 0:
            for spec in runs_spec:
                r = p.add_run()
                r.text = ""
                _apply_font(r, spec["font"])
        else:
            base, rem = divmod(total, n)
            sizes = [base + (1 if k < rem else 0) for k in range(n)]
            pos = 0
            for k, spec in enumerate(runs_spec):
                seg = line[pos:pos + sizes[k]]
                pos += sizes[k]
                r = p.add_run()
                r.text = seg
                _apply_font(r, spec["font"])

def shape_by_id(slide, shape_id):
    sid = int(shape_id)
    for sh in slide.shapes:
        val = getattr(sh, "shape_id", None) or getattr(getattr(sh, "_element", None), "cNvPr", None)
        val = int(getattr(val, "id", val)) if val is not None else None
        if val == sid:
            return sh
    return None


def ensure_slot_textbox(slide, slot):
    slot = (slot or "").lower()

    def _get(ph_name):
        return getattr(PP_PLACEHOLDER, ph_name, None)

    TITLE = _get("TITLE")
    CENTER_TITLE = _get("CENTER_TITLE")
    SUBTITLE = _get("SUBTITLE")
    BODY = _get("BODY")
    CONTENT = _get("CONTENT")
    OBJECT = _get("OBJECT")

    title_types = {t for t in (TITLE, CENTER_TITLE, SUBTITLE) if t is not None}
    body_types  = {t for t in (BODY, CONTENT, OBJECT) if t is not None}

    def find_placeholder(accepted_types):
        for sh in slide.shapes:
            if not getattr(sh, "is_placeholder", False):
                continue
            pf = getattr(sh, "placeholder_format", None)
            if not pf:
                continue
            try:
                if pf.type in accepted_types:
                    return sh
            except Exception:
                pass
        return None

    if slot == "title":
        ph = find_placeholder(title_types)
        if ph:
            return ph

    if slot == "body":
        ph = find_placeholder(body_types)
        if ph:
            return ph

    from pptx.util import Inches
    if slot == "title":
        return slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    if slot == "body":
        return slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    return slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))

def _layout_has(layout, want_title=False, want_body=False):
    has_title = has_body = False
    for ph in getattr(layout, "placeholders", []):
        pf = getattr(ph, "placeholder_format", None)
        t = getattr(pf, "type", None) if pf else None
        if t in (getattr(PP_PLACEHOLDER, "TITLE", None),
                 getattr(PP_PLACEHOLDER, "CENTER_TITLE", None),
                 getattr(PP_PLACEHOLDER, "SUBTITLE", None)):
            has_title = True
        if t in (getattr(PP_PLACEHOLDER, "BODY", None),
                 getattr(PP_PLACEHOLDER, "CONTENT", None),
                 getattr(PP_PLACEHOLDER, "OBJECT", None)):
            has_body = True
    return (not want_title or has_title) and (not want_body or has_body)

def _pick_layout_for_slots(prs, anchor_slide, needs_title, needs_body):
    if anchor_slide and _layout_has(anchor_slide.slide_layout, needs_title, needs_body):
        return anchor_slide.slide_layout
    for layout in prs.slide_layouts:
        if _layout_has(layout, needs_title, needs_body):
            return layout
    return anchor_slide.slide_layout if anchor_slide else prs.slide_layouts[-1]

def _collect_needs(edit_items):
    needs = {}
    for tgt, _ in edit_items:
        if not isinstance(tgt, str):
            continue
        m = re.match(r"^(n\d+):slot:(title|body)$", tgt.strip(), flags=re.I)
        if m:
            ref, slot = m.group(1), m.group(2).lower()
            needs.setdefault(ref, {"title": False, "body": False})
            needs[ref][slot] = True
    return needs

def _body_placeholder_bounds(slide):
    """Return (left, top, width, height) for the body/content area if possible, else None."""
    try:
        for shp in slide.shapes:
            phf = getattr(shp, "placeholder_format", None)
            if phf is not None:
                # BODY placeholder is the content region on most layouts
                if str(getattr(phf, "type", "")).endswith("BODY"):
                    return shp.left, shp.top, shp.width, shp.height
    except Exception:
        pass
    return None

def _add_table_from_matrix(slide, matrix):
    """
    Create a table on the slide sized to the matrix (rows x cols) and fill it.
    The table is placed over the body placeholder bounds if available,
    else within 1-inch margins.
    Returns the created table shape.
    """
    if not isinstance(matrix, (list, tuple)) or not matrix or not isinstance(matrix[0], (list, tuple)):
        return None

    rows = len(matrix)
    cols = len(matrix[0])

    # determine placement rectangle
    rect = _body_placeholder_bounds(slide)
    if rect:
        left, top, width, height = rect
    else:
        # safe default margins
        left = Inches(1)
        top = Inches(1.2)
        # try to use slide size when available
        try:
            prs = slide.part.presentation
            width = prs.slide_width - Inches(2)
            height = prs.slide_height - Inches(2.2)
        except Exception:
            width = Inches(8)
            height = Inches(4.5)

    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = tbl_shape.table

    # fill cells
    for r in range(rows):
        for c in range(cols):
            try:
                table.cell(r, c).text = "" if matrix[r][c] is None else str(matrix[r][c])
            except Exception:
                pass

    return tbl_shape


def _resolve_donor_simple(order, slides_by_id, anchor_id, kind):
    """
    kind: 'insert_after' or 'insert_before'
    Rules:
      insert_after(anchor):
        - if anchor is first -> donor = next slide if exists else anchor
        - else               -> donor = anchor
      insert_before(anchor):
        - if anchor is last  -> donor = previous slide if exists else anchor
        - else               -> donor = anchor
    """
    if not order:
        return None
    if anchor_id not in order:
        # anchor not found
        return slides_by_id.get(order[1]) if len(order) > 1 else slides_by_id.get(order[0])

    pos = order.index(anchor_id)
    last_idx = len(order) - 1

    if kind == "insert_after":
        if pos == 0:
            # after first slide
            return slides_by_id.get(order[pos + 1]) if pos + 1 <= last_idx else slides_by_id.get(anchor_id)
        else:
            return slides_by_id.get(anchor_id)

    # insert_before
    if pos == last_idx:
        # before last slide
        return slides_by_id.get(order[pos - 1]) if pos - 1 >= 0 else slides_by_id.get(anchor_id)
    else:
        return slides_by_id.get(anchor_id)

def _set_table_from_matrix(shape, data):
    # data = list[list[Any]]; trims to current table size
    tbl = shape.table
    max_r = len(tbl.rows)
    max_c = len(tbl.columns)
    for r, row_vals in enumerate(data):
        if r >= max_r:
            break
        for c, val in enumerate(row_vals):
            if c >= max_c:
                break
            tbl.cell(r, c).text = ""  # clear
            tbl.cell(r, c).text = "" if val is None else str(val)


@mcp.tool()
async def edit_document(
    file_id: str,
    file_name: str,
    operations: list[dict],
    ctx: Context[ServerSession, None] = None
) -> dict:
    """
    Edit a document (DOCX, XLSX, PPTX) with simple, natural operations.
    
    **ALWAYS call `files_full_context_document()` first to see document structure!**
    
    **Simple 3-Field Format:**
```javascript
    operations: [
      {
        action: "append",     // What to do
        content: "text",      // What to add (string or structured)
        target: "pid:1"       // Where (optional, depends on action)
      }
    ]
```
    
    **DOCX Actions:**
    
    **1. Append text to end** (most common):
```javascript
    {action: "append", content: "Simple paragraph"}
    
    // Multiple paragraphs:
    {action: "append", content: ["Para 1", "Para 2", "Para 3"]}
    
    // Styled content (headings, etc):
    {action: "append", content: [
      {style: "Heading 1", text: "Title"},
      {style: "Heading 2", text: "Subtitle"},
      {style: "Normal", text: "Body text"}
    ]}
```
    
    **2. Replace paragraph:**
```javascript
    {action: "replace", target: "pid:3", content: "New text"}
    
    // Or find by content:
    {action: "replace", target: {find: "old text"}, content: "New text"}
```
    
    **3. Insert after/before:**
```javascript
    {action: "insert_after", target: "pid:2", content: "New paragraph"}
    {action: "insert_before", target: "pid:1", content: "New paragraph"}
```
    
    **4. Delete:**
```javascript
    {action: "delete", target: "pid:5"}
```
    
    **Supported Styles (DOCX):**
    - `"Heading 1"`, `"Heading 2"`, `"Heading 3"`, etc.
    - `"Title"`
    - `"Normal"` (default)
    - `"Quote"`, `"Intense Quote"`
    
    **XLSX Actions:**
```javascript
    {action: "replace", target: "A1", content: "value"}
    {action: "replace", target: "B5", content: 42}
```
    
    **Returns:**
```json
    {
      "url": "https://files.example.com/edited.docx",
      "operations_applied": 3,
      "changes": ["✓ Appended 5 paragraphs", "✓ Replaced pid:2"]
    }
```
    """
    
    # Auto-deserialize stringified operations (some models send JSON strings)
    if isinstance(operations, str):
        try:
            operations = json.loads(operations)
        except (json.JSONDecodeError, TypeError):
            return {
                "error": "operations must be a list (received unparseable string)",
                "example": '[{"action": "append", "content": "text"}]'
            }
    if not operations or not isinstance(operations, list):
        return {
            "error": "operations must be a non-empty list",
            "example": '[{"action": "append", "content": "text"}]'
        }
    
    # === AUTH & DOWNLOAD ===
    bearer_token = _try_get_auth_header(ctx)
    user_token = bearer_token if bearer_token else None
    
    user_file = download_file(file_id=file_id, token=user_token)
    if isinstance(user_file, dict) and "error" in user_file:
        return user_file
    
    file_extension = os.path.splitext(file_name)[1].lower()
    file_type = file_extension.lstrip('.')
    
    temp_folder = f"/app/temp/{uuid.uuid4()}"
    os.makedirs(temp_folder, exist_ok=True)
    
    try:
        changes_log = []
        
        # =====================================================================
        # DOCX EDITING
        # =====================================================================
        if file_type == "docx":
            doc = Document(user_file)
            
            # Build paragraph registry
            paragraphs = []
            for idx, para in enumerate(doc.paragraphs, 1):
                if para.text.strip():
                    paragraphs.append({
                        "id": f"pid:{idx}",
                        "index": idx,
                        "element": para,
                        "text": para.text,
                        "style": para.style.name if para.style else "Normal"
                    })
            
            # Helper: Find paragraph
            def find_paragraph(target):
                if isinstance(target, str):
                    for p in paragraphs:
                        if p["id"] == target:
                            return p
                elif isinstance(target, dict) and "find" in target:
                    search_text = target["find"].lower()
                    for p in paragraphs:
                        if search_text in p["text"].lower():
                            return p
                return None
            
            # Helper: Add styled content
            def add_content(content, reference_para=None):
                """Add content with proper styling. Returns list of new paragraphs."""
                new_paras = []
                
                # Determine what kind of content we have
                if isinstance(content, str):
                    # Simple string - single paragraph
                    items = [{"style": "Normal", "text": content}]
                elif isinstance(content, list):
                    if all(isinstance(x, str) for x in content):
                        # List of strings - multiple normal paragraphs
                        items = [{"style": "Normal", "text": t} for t in content]
                    elif all(isinstance(x, dict) for x in content):
                        # List of styled objects
                        items = content
                    else:
                        # Mixed - treat as strings
                        items = [{"style": "Normal", "text": str(t)} for t in content]
                else:
                    items = [{"style": "Normal", "text": str(content)}]
                
                # Map block type names → Word style names
                _TYPE_TO_STYLE = {
                    "title": "Title",
                    "heading": "Heading 1",
                    "heading1": "Heading 1",
                    "heading2": "Heading 2",
                    "heading3": "Heading 3",
                    "subheading": "Heading 2",
                    "subtitle": "Subtitle",
                    "paragraph": "Normal",
                    "normal": "Normal",
                    "body": "Normal",
                    "list": "List Bullet",
                    "bullet": "List Bullet",
                    "quote": "Quote",
                    "blockquote": "Quote",
                }

                for item in items:
                    if isinstance(item, dict):
                        text = item.get("text", str(item))
                        # Priority: explicit style > mapped type > Normal
                        style = item.get("style")
                        if not style:
                            block_type = item.get("type", "").lower()
                            style = _TYPE_TO_STYLE.get(block_type, "Normal")
                    else:
                        text = str(item)
                        style = "Normal"
                    
                    # Create new paragraph
                    new_para = doc.add_paragraph(text)
                    
                    # Apply style
                    try:
                        new_para.style = style
                    except:
                        # If style doesn't exist, use Normal
                        new_para.style = "Normal"
                        changes_log.append(f"⚠️ Style '{style}' not found, using Normal")
                    
                    # Copy formatting from reference if provided
                    if reference_para and not isinstance(item, dict):
                        try:
                            new_para.style = reference_para.style
                        except:
                            pass
                    
                    new_paras.append(new_para)
                
                return new_paras
            
            # Process operations
            for op_idx, op in enumerate(operations):
                if not isinstance(op, dict):
                    changes_log.append(f"⚠️ Skipped operation {op_idx+1}: not a dict")
                    continue
                
                # Normalize key aliases — models hallucinate varying key names
                # Action key: action, op, type, operation
                if "action" not in op:
                    for alt in ("op", "type", "operation", "command", "method"):
                        if alt in op:
                            op["action"] = op.pop(alt)
                            break
                # Target key: target, target_id, target_pid, id, pid, ref
                if "target" not in op:
                    for alt in ("target_id", "target_pid", "id", "pid", "ref", "anchor", "after", "before"):
                        if alt in op:
                            op["target"] = op.pop(alt)
                            break
                # Content key: content, text, value, data, body
                if "content" not in op:
                    for alt in ("text", "value", "data", "body", "new_content", "new_text"):
                        if alt in op and alt != "text":
                            op["content"] = op.pop(alt)
                            break
                        elif alt == "text" and alt in op and isinstance(op[alt], str) and "action" in op:
                            # Only steal "text" if there's already an action (avoid grabbing text from content blocks)
                            op["content"] = op.pop(alt)
                            break
                
                if "action" not in op:
                    changes_log.append(f"⚠️ Skipped operation {op_idx+1}: missing 'action'")
                    continue
                
                action = op.get("action", "").lower().strip()
                target = op.get("target")
                content = op.get("content")
                
                # Normalize hallucinated action names → valid actions
                _ACTION_ALIASES = {
                    "append_paragraph": "append",
                    "add": "append",
                    "add_paragraph": "append",
                    "add_section": "append",
                    "append_section": "append",
                    "append_text": "append",
                    "replace_paragraph": "replace",
                    "update": "replace",
                    "edit": "replace",
                    "modify": "replace",
                    "remove": "delete",
                    "delete_paragraph": "delete",
                    "insert": "insert_after",
                    "insert_paragraph": "insert_after",
                    "add_after": "insert_after",
                    "add_before": "insert_before",
                }
                action = _ACTION_ALIASES.get(action, action)
                
                # APPEND - Add to end
                if action == "append":
                    if content is None:
                        changes_log.append(f"⚠️ Skipped append: missing 'content'")
                        continue
                    
                    new_paras = add_content(content)
                    count = len(new_paras)
                    changes_log.append(f"✓ Appended {count} paragraph{'s' if count != 1 else ''}")
                    continue
                
                # REPLACE
                if action == "replace":
                    para_info = find_paragraph(target)
                    if not para_info:
                        available = [p["id"] for p in paragraphs[:5]]
                        return {
                            "error": f"Target '{target}' not found",
                            "available_targets": available + ["..."] if len(paragraphs) > 5 else available
                        }
                    
                    para = para_info["element"]
                    para.clear()
                    para.add_run(str(content))
                    changes_log.append(f"✓ Replaced {para_info['id']}")
                    continue
                
                # INSERT_AFTER
                if action == "insert_after":
                    para_info = find_paragraph(target)
                    if not para_info:
                        return {
                            "error": f"Target '{target}' not found",
                            "available_targets": [p["id"] for p in paragraphs[:5]]
                        }
                    
                    anchor = para_info["element"]
                    new_paras = add_content(content, reference_para=anchor)
                    
                    # Insert into XML tree
                    anchor_element = anchor._element
                    parent = anchor_element.getparent()
                    insert_idx = parent.index(anchor_element) + 1
                    
                    for new_para in reversed(new_paras):  # Reverse to maintain order
                        parent.insert(insert_idx, new_para._element)
                    
                    changes_log.append(f"✓ Inserted {len(new_paras)} paragraph(s) after {para_info['id']}")
                    continue
                
                # INSERT_BEFORE
                if action == "insert_before":
                    para_info = find_paragraph(target)
                    if not para_info:
                        return {
                            "error": f"Target '{target}' not found",
                            "available_targets": [p["id"] for p in paragraphs[:5]]
                        }
                    
                    anchor = para_info["element"]
                    new_paras = add_content(content, reference_para=anchor)
                    
                    anchor_element = anchor._element
                    parent = anchor_element.getparent()
                    insert_idx = parent.index(anchor_element)
                    
                    for idx, new_para in enumerate(new_paras):
                        parent.insert(insert_idx + idx, new_para._element)
                    
                    changes_log.append(f"✓ Inserted {len(new_paras)} paragraph(s) before {para_info['id']}")
                    continue
                
                # DELETE
                if action == "delete":
                    para_info = find_paragraph(target)
                    if not para_info:
                        return {
                            "error": f"Target '{target}' not found",
                            "available_targets": [p["id"] for p in paragraphs[:5]]
                        }
                    
                    para = para_info["element"]
                    parent = para._element.getparent()
                    parent.remove(para._element)
                    changes_log.append(f"✓ Deleted {para_info['id']}")
                    continue
                
                changes_log.append(f"⚠️ Unknown action: '{action}'")
            
            # Save
            edited_path = os.path.join(temp_folder, f"{os.path.splitext(file_name)[0]}_edited.docx")
            doc.save(edited_path)
        
        # =====================================================================
        # XLSX EDITING
        # =====================================================================
        elif file_type == "xlsx":
            wb = load_workbook(user_file)
            ws = wb.active
            
            for op in operations:
                if not isinstance(op, dict):
                    continue
                
                # Normalize key aliases — models hallucinate varying key names
                # Action key: action, op, type, operation
                if "action" not in op:
                    for alt in ("op", "type", "operation", "command", "method"):
                        if alt in op:
                            op["action"] = op.pop(alt)
                            break
                # Target key: target, target_id, cell, id
                if "target" not in op:
                    for alt in ("target_id", "cell", "id", "ref", "address"):
                        if alt in op:
                            op["target"] = op.pop(alt)
                            break
                # Content key: content, value, data, text
                if "content" not in op:
                    for alt in ("value", "data", "text", "body", "new_value"):
                        if alt in op:
                            op["content"] = op.pop(alt)
                            break
                
                if "action" not in op:
                    continue
                
                action = op.get("action", "").lower().strip()
                target = op.get("target", "A1")
                content = op.get("content")
                
                # Normalize hallucinated action names
                _XLSX_ALIASES = {
                    "set": "replace", "update": "replace", "edit": "replace",
                    "modify": "replace", "change": "replace", "write": "replace",
                    "set_cell": "replace", "update_cell": "replace",
                    "replace_cell": "replace",
                }
                action = _XLSX_ALIASES.get(action, action)
                
                if action == "replace":
                    try:
                        ws[target] = content
                        changes_log.append(f"✓ Set {target} = {content}")
                    except Exception as e:
                        changes_log.append(f"⚠️ Failed to set {target}: {e}")
            
            edited_path = os.path.join(temp_folder, f"{os.path.splitext(file_name)[0]}_edited.xlsx")
            wb.save(edited_path)
        
        # =====================================================================
        # PPTX EDITING  
        # =====================================================================
        elif file_type == "pptx":
            prs = Presentation(user_file)
            slides_by_id = {int(s.slide_id): s for s in prs.slides}
            
            for op in operations:
                if not isinstance(op, dict) or "action" not in op:
                    continue
                
                action = op.get("action", "").lower().strip()
                target = op.get("target")
                content = op.get("content")
                
                # Normalize hallucinated action names
                _PPTX_ALIASES = {
                    "update": "replace", "edit": "replace", "modify": "replace",
                    "change": "replace", "set": "replace",
                    "replace_text": "replace", "update_text": "replace",
                }
                action = _PPTX_ALIASES.get(action, action)
                
                if action == "replace" and target:
                    m = re.match(r"^sid:(\d+)/shid:(\d+)$", str(target))
                    if m:
                        slide_id = int(m.group(1))
                        shape_id = int(m.group(2))
                        slide = slides_by_id.get(slide_id)
                        if slide:
                            shape = shape_by_id(slide, shape_id)
                            if shape and hasattr(shape, "text_frame"):
                                shape.text_frame.text = str(content)
                                changes_log.append(f"✓ Updated {target}")
            
            edited_path = os.path.join(temp_folder, f"{os.path.splitext(file_name)[0]}_edited.pptx")
            prs.save(edited_path)
        
        else:
            raise Exception(f"Unsupported file type: {file_type}")
        
        # Upload
        response = upload_file(
            file_path=edited_path,
            filename=f"{os.path.splitext(file_name)[0]}_edited",
            file_type=file_type,
            token=user_token
        )
        
        if isinstance(response, dict):
            response["operations_applied"] = len([c for c in changes_log if c.startswith("✓")])
            response["changes"] = changes_log
        
        shutil.rmtree(temp_folder, ignore_errors=True)
        return response
    
    except Exception as e:
        shutil.rmtree(temp_folder, ignore_errors=True)
        return {
            "error": str(e),
            "changes_attempted": changes_log
        }


# ============================================================================
# UPDATED download_file FUNCTION (also needed)
# ============================================================================

def download_file(file_id: str, token: str = None) -> BytesIO | dict:
    """
    Download a file from Open WebUI (disk fast-path → HTTP fallback).
    
    Returns:
        BytesIO on success
        dict with "error" key on failure
    """
    base_dir = "/openwebui_uploads"
    
    # === DISK FAST-PATH (no auth needed) ===
    if os.path.exists(base_dir):
        try:
            # A) Flat layout: /openwebui_uploads/<file_id>_<original_filename>
            import glob
            flat_matches = glob.glob(os.path.join(base_dir, f"{file_id}_*"))
            flat_matches = [p for p in flat_matches if os.path.isfile(p)]
            if flat_matches:
                file_path = sorted(flat_matches)[0]
                logging.info(f"[disk-first] Reading: {file_path}")
                with open(file_path, "rb") as f:
                    return BytesIO(f.read())

            # B) Legacy dir layout: /openwebui_uploads/<file_id>/<somefile>
            id_dir = os.path.join(base_dir, file_id)
            if os.path.isdir(id_dir):
                files = [f for f in os.listdir(id_dir) if os.path.isfile(os.path.join(id_dir, f))]
                if files:
                    file_path = os.path.join(id_dir, sorted(files)[0])
                    logging.info(f"[disk-first] Reading (legacy): {file_path}")
                    with open(file_path, "rb") as f:
                        return BytesIO(f.read())
                        
        except Exception as e:
            logging.warning(f"[disk-first] Read failed: {e}")

    # === HTTP FALLBACK (only if token provided) ===
    if not token:
        return {"error": {"message": f"File {file_id} not found on disk and no auth token for HTTP fallback"}}
    
    logging.info(f"[http-fallback] Downloading {file_id}")
    url = f"{URL}/api/v1/files/{file_id}/content"
    
    auth_header = token if token.startswith("Bearer ") else f"Bearer {token}"
    headers = {
        "Authorization": auth_header,
        "Accept": "application/json"
    }
    
    try:
        response = get(url, headers=headers, timeout=30)
        if response.status_code != 200:
            return {"error": {"message": f"HTTP {response.status_code} downloading file {file_id}"}}
        return BytesIO(response.content)
    except Exception as e:
        return {"error": {"message": f"HTTP download failed: {e}"}}

def _get_pptx_namespaces():
    """Returns XML namespaces for PowerPoint"""
    return {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p15': 'http://schemas.microsoft.com/office/powerpoint/2012/main',
        'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main'
    }

def _add_native_pptx_comment_zip(pptx_path, slide_num, comment_text, author_id, x=100, y=100):
    """
    Add a native PowerPoint comment by directly manipulating the ZIP file.
        Args:
        pptx_path: Path to the PPTX file
        slide_num: Slide number (1-based)
        comment_text: Comment text
        author_id: Author ID
        x: X position in EMU (not pixels!)
        y: Y position in EMU (not pixels!)
    """
    namespaces = _get_pptx_namespaces()
    
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            zf.extractall(temp_path)
        
        authors_file = temp_path / 'ppt' / 'commentAuthors.xml'
        if authors_file.exists():
            root = etree.parse(str(authors_file)).getroot()
            found = False
            for author in root.findall('.//p:cmAuthor', namespaces):
                if author.get('name') == 'AI Reviewer':
                    author_id = int(author.get('id'))
                    found = True
                    break
            
            if not found:
                existing_ids = [int(a.get('id')) for a in root.findall('.//p:cmAuthor', namespaces)]
                author_id = max(existing_ids) + 1 if existing_ids else 0
                author = etree.SubElement(root, f'{{{namespaces["p"]}}}cmAuthor')
                author.set('id', str(author_id))
                author.set('name', 'AI Reviewer')
                author.set('initials', 'AI')
                author.set('lastIdx', '1')
                author.set('clrIdx', str(author_id % 8))
        else:
            authors_file.parent.mkdir(parents=True, exist_ok=True)
            root = etree.Element(
                f'{{{namespaces["p"]}}}cmAuthorLst',
                nsmap={k: v for k, v in namespaces.items() if k in ['p']}
            )
            author = etree.SubElement(root, f'{{{namespaces["p"]}}}cmAuthor')
            author.set('id', str(author_id))
            author.set('name', 'AI Reviewer')
            author.set('initials', 'AI')
            author.set('lastIdx', '1')
            author.set('clrIdx', '0')
            
            rels_file = temp_path / 'ppt' / '_rels' / 'presentation.xml.rels'
            if rels_file.exists():
                rels_root = etree.parse(str(rels_file)).getroot()
                existing_ids = [int(rel.get('Id')[3:]) for rel in rels_root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')]
                next_rid = max(existing_ids) + 1 if existing_ids else 1
                
                rel = etree.SubElement(rels_root, '{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')
                rel.set('Id', f'rId{next_rid}')
                rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/commentAuthors')
                rel.set('Target', 'commentAuthors.xml')
                
                with open(rels_file, 'wb') as f:
                    f.write(etree.tostring(rels_root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        with open(authors_file, 'wb') as f:
            f.write(etree.tostring(root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        comments_dir = temp_path / 'ppt' / 'comments'
        comments_dir.mkdir(parents=True, exist_ok=True)
        comment_file = comments_dir / f'comment{slide_num}.xml'
        
        if comment_file.exists():
            comments_root = etree.parse(str(comment_file)).getroot()
        else:
            comments_root = etree.Element(
                f'{{{namespaces["p"]}}}cmLst',
                nsmap={k: v for k, v in namespaces.items() if k in ['p']}
            )
            
            slide_rels_file = temp_path / 'ppt' / 'slides' / '_rels' / f'slide{slide_num}.xml.rels'
            if slide_rels_file.exists():
                slide_rels_root = etree.parse(str(slide_rels_file)).getroot()
            else:
                slide_rels_file.parent.mkdir(parents=True, exist_ok=True)
                slide_rels_root = etree.Element(
                    '{http://schemas.openxmlformats.org/package/2006/relationships}Relationships'
                )
            
            existing_ids = [int(rel.get('Id')[3:]) for rel in slide_rels_root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')]
            next_rid = max(existing_ids) + 1 if existing_ids else 1
            
            rel = etree.SubElement(slide_rels_root, '{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')
            rel.set('Id', f'rId{next_rid}')
            rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments')
            rel.set('Target', f'../comments/comment{slide_num}.xml')
            
            with open(slide_rels_file, 'wb') as f:
                f.write(etree.tostring(slide_rels_root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        existing_ids = [int(c.get('idx')) for c in comments_root.findall('.//p:cm', namespaces)]
        next_id = max(existing_ids) + 1 if existing_ids else 1
        
        comment = etree.SubElement(comments_root, f'{{{namespaces["p"]}}}cm')
        comment.set('authorId', str(author_id))
        comment.set('dt', datetime.datetime.now().isoformat())
        comment.set('idx', str(next_id))
        
        pos = etree.SubElement(comment, f'{{{namespaces["p"]}}}pos')
        pos.set('x', str(int(x)))
        pos.set('y', str(int(y)))
        
        text_elem = etree.SubElement(comment, f'{{{namespaces["p"]}}}text')
        text_elem.text = comment_text
        
        with open(comment_file, 'wb') as f:
            f.write(etree.tostring(comments_root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        content_types_file = temp_path / '[Content_Types].xml'
        if content_types_file.exists():
            ct_root = etree.parse(str(content_types_file)).getroot()
            ns = {'ct': 'http://schemas.openxmlformats.org/package/2006/content-types'}
            
            has_authors = False
            has_comments = False
            
            for override in ct_root.findall('.//ct:Override', ns):
                if override.get('PartName') == '/ppt/commentAuthors.xml':
                    has_authors = True
                if override.get('PartName') == f'/ppt/comments/comment{slide_num}.xml':
                    has_comments = True
            
            if not has_authors:
                override = etree.SubElement(ct_root, '{http://schemas.openxmlformats.org/package/2006/content-types}Override')
                override.set('PartName', '/ppt/commentAuthors.xml')
                override.set('ContentType', 'application/vnd.openxmlformats-officedocument.presentationml.commentAuthors+xml')
            
            if not has_comments:
                override = etree.SubElement(ct_root, '{http://schemas.openxmlformats.org/package/2006/content-types}Override')
                override.set('PartName', f'/ppt/comments/comment{slide_num}.xml')
                override.set('ContentType', 'application/vnd.openxmlformats-officedocument.presentationml.comments+xml')
            
            with open(content_types_file, 'wb') as f:
                f.write(etree.tostring(ct_root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for file_path in temp_path.rglob('*'):
                if file_path.is_file():
                    arcname = str(file_path.relative_to(temp_path))
                    zf.write(file_path, arcname)
        
        log.debug(f"Native comment added to slide {slide_num} with idx={next_id}")

@mcp.tool(
    name="review_document",
    title="Review and comment on various document types",
    description="Review an existing document of various types (docx, xlsx, pptx), perform corrections and add comments."
)
async def review_document(
    file_id: str,
    file_name: str,
    review_comments: List[ReviewComment],
    ctx: Context[ServerSession, None]
) -> dict:
    """
    Generic document review function that works with different document types.
    File type is automatically detected from the file extension.
    Returns a markdown hyperlink for downloading the reviewed document.
    
    For Excel files (.xlsx):
    - The index must be a cell reference (e.g. "A1", "B3", "C10")
    - These correspond to the "index" key returned by the full_context_document() function
    - Never use integer values for Excel cells
    
    For Word files (.docx):
    - The index should be a paragraph ID in the format "pid:<para_xml_id>"
    - These correspond to the "id_key" field returned by the full_context_document() function
    
    For PowerPoint files (.pptx):
    - The index should be a slide ID in the format "sid:<slide_id>"
    - These correspond to the "id_key" field returned by the full_context_document() function
    """
    temp_folder = f"/app/temp/{uuid.uuid4()}"
    os.makedirs(temp_folder, exist_ok=True)

    try:
        bearer_token = ctx.request_context.request.headers.get("authorization")
        logging.info(f"Recieved authorization header!")
        user_token=bearer_token
    except:
        logging.error(f"Error retrieving authorization header use admin fallback")
        user_token=TOKEN
    try:
        user_file = download_file(file_id=file_id, token=user_token)
        if isinstance(user_file, dict) and "error" in user_file:
            return json.dumps(user_file, indent=4, ensure_ascii=False)

        file_extension = os.path.splitext(file_name)[1].lower()
        file_type = file_extension.lstrip('.')

        reviewed_path = None
        response = None

        # Normalize to list of objects {index, comment}
        norm_comments: List[ReviewComment] = []
        try:
            if isinstance(review_comments, list):
                if len(review_comments) == 0:
                    norm_comments = []
                elif isinstance(review_comments[0], dict):
                    for item in review_comments:
                        if isinstance(item, dict) and "index" in item and "comment" in item:
                            norm_comments.append({"index": item["index"], "comment": str(item["comment"])})
                else:
                    # Legacy: [[index, comment], ...] or tuples
                    for item in review_comments:
                        if isinstance(item, (list, tuple)) and len(item) >= 2:
                            norm_comments.append({"index": item[0], "comment": str(item[1])})
        except Exception:
            pass

        if file_type == "docx":
            try:
                doc = Document(user_file)
                paragraphs = list(doc.paragraphs)
                para_by_xml_id = {}
                para_id_counter = 1
                
                for para in doc.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    para_by_xml_id[para_id_counter] = para
                    para_id_counter += 1

                for rc in norm_comments:
                    index = rc["index"]
                    comment_text = rc["comment"]
                    if isinstance(index, int) and 0 <= index < len(paragraphs):
                        para = paragraphs[index]
                        if para.runs:
                            try:
                                doc.add_comment(
                                    runs=[para.runs[0]],
                                    text=comment_text,
                                    author="AI Reviewer",
                                    initials="AI"
                                )
                            except Exception:
                                para.add_run(f"  [AI Comment: {comment_text}]")
                    elif isinstance(index, str) and index.startswith("pid:"):
                        try:
                            para_xml_id = int(index.split(":")[1])
                            para = para_by_xml_id.get(para_xml_id)
                            if para and para.runs:
                                try:
                                    doc.add_comment(
                                        runs=[para.runs[0]],
                                        text=comment_text,
                                        author="AI Reviewer",
                                        initials="AI"
                                    )
                                except Exception:
                                    para.add_run(f"  [AI Comment: {comment_text}]")
                        except Exception:
                            if isinstance(index, int) and 0 <= index < len(paragraphs):
                                para = paragraphs[index]
                                if para.runs:
                                    try:
                                        doc.add_comment(
                                            runs=[para.runs[0]],
                                            text=comment_text,
                                            author="AI Reviewer",
                                            initials="AI"
                                        )
                                    except Exception:
                                        para.add_run(f"  [AI Comment: {comment_text}]")
                reviewed_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_reviewed.docx"
                )
                doc.save(reviewed_path)
                response = upload_file(
                    file_path=reviewed_path,
                    filename=f"{os.path.splitext(file_name)[0]}_reviewed",
                    file_type="docx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error during DOCX revision: {e}")

        elif file_type == "xlsx":
            try:
                wb = load_workbook(user_file)
                ws = wb.active

                for rc in norm_comments:
                    index = rc["index"]
                    comment_text = rc["comment"]
                    try:
                        if isinstance(index, str) and re.match(r"^[A-Z]+[0-9]+$", index.strip().upper()):
                            cell_ref = index.strip().upper()
                        elif isinstance(index, int):
                            cell_ref = f"A{index+1}"
                        else:
                            cell_ref = "A1"

                        cell = ws[cell_ref]
                        add_auto_sized_review_comment(cell, comment_text, author="AI Reviewer")

                    except Exception:
                        fallback_cell = ws["A1"]
                        add_auto_sized_review_comment(fallback_cell, comment_text, author="AI Reviewer")

                reviewed_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_reviewed.xlsx"
                )
                wb.save(reviewed_path)
                response = upload_file(
                    file_path=reviewed_path,
                    filename=f"{os.path.splitext(file_name)[0]}_reviewed",
                    file_type="xlsx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error: {e}")

        elif file_type == "pptx":
            try:
                temp_pptx = os.path.join(temp_folder, "temp_input.pptx")
                with open(temp_pptx, 'wb') as f:
                    f.write(user_file.read())
                
                prs = Presentation(temp_pptx)
                slides_by_id = {int(s.slide_id): s for s in prs.slides}
                
                comments_by_slide = {}
                
                for rc in norm_comments:
                    index = rc["index"]
                    comment_text = rc["comment"]
                    slide_num = None
                    slide_id = None
                    
                    if isinstance(index, int) and 0 <= index < len(prs.slides):
                        slide_num = index + 1
                        slide_id = list(slides_by_id.keys())[index]
                    elif isinstance(index, str):
                        if index.startswith("sid:") and "/shid:" in index:
                            try:
                                slide_id = int(index.split("/")[0].replace("sid:", ""))
                                if slide_id in slides_by_id:
                                    slide_num = list(slides_by_id.keys()).index(slide_id) + 1
                            except Exception as e:
                                log.warning(f"Failed to parse shape ID: {e}")
                        elif index.startswith("sid:"):
                            try:
                                slide_id = int(index.split(":")[1])
                                if slide_id in slides_by_id:
                                    slide_num = list(slides_by_id.keys()).index(slide_id) + 1
                            except Exception as e:
                                log.warning(f"Failed to parse slide ID: {e}")
                    
                    if slide_num and slide_id:
                        if slide_num not in comments_by_slide:
                            comments_by_slide[slide_num] = []
                        
                        shape_info = ""
                        if "/shid:" in str(index):
                            try:
                                shape_id = int(str(index).split("/shid:")[1])
                                shape_info = f"[Shape {shape_id}] "
                            except:
                                pass
                        
                        comments_by_slide[slide_num].append(f"{shape_info}{comment_text}")
                comment_offset = 0              
                for slide_num, comments in comments_by_slide.items():
                    comment_start_x = 5000
                    comment_start_y = 1000
                    comment_spacing_y = 1500
                    
                    for idx, comment in enumerate(comments):
                        try:
                            y_position = comment_start_y + (idx * comment_spacing_y)
                            
                            _add_native_pptx_comment_zip(
                                pptx_path=temp_pptx,
                                slide_num=slide_num,
                                comment_text=f"• {comment}",
                                author_id=0,
                                x=comment_start_x,
                                y=y_position
                            )
                            log.debug(f"Native PowerPoint comment added to slide {slide_num} at position x={comment_start_x}, y={y_position}")
                        except Exception as e:
                            log.warning(f"Failed to add native comment to slide {slide_num}: {e}", exc_info=True)
                            prs_fallback = Presentation(temp_pptx)
                            slide = prs_fallback.slides[slide_num - 1]
                            left = top = Inches(0.2)
                            width = Inches(4)
                            height = Inches(1)
                            textbox = slide.shapes.add_textbox(left, top, width, height)
                            text_frame = textbox.text_frame
                            p = text_frame.add_paragraph()
                            p.text = f"AI Reviewer: {comment}"
                            p.font.size = PptPt(10)
                            prs_fallback.save(temp_pptx)

                reviewed_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_reviewed.pptx"
                )
                shutil.copy(temp_pptx, reviewed_path)
                
                response = upload_file(
                    file_path=reviewed_path,
                    filename=f"{os.path.splitext(file_name)[0]}_reviewed",
                    file_type="pptx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error when revising PPTX: {e}")

        else:
            raise Exception(f"File type not supported : {file_type}")

        shutil.rmtree(temp_folder, ignore_errors=True)

        return response

    except Exception as e:
        shutil.rmtree(temp_folder, ignore_errors=True)
        return json.dumps(
            {"error": {"message": str(e)}},
            indent=4,
            ensure_ascii=False
        )
# ─── PATCH: Replace the create_file function in file_export_mcp.py ───
# Replaces everything from @mcp.tool() above create_file to return {"url": result["url"]}
#
# Changes:
#   1. B5 FIX: Unwrap structured content blocks from OpenRouter/ChatGPT models
#      Models send: content=[{"type":"raw","format":"json","text":"...actual..."}]
#      instead of: content="...actual..."
#   2. Stringified JSON fallback: content='[{"type":"paragraph",...}]' as string → parsed list
#   3. Better logging for debugging model content format issues
#   4. Consolidated content normalization before format routing

# ── Text-based formats: content should be a raw string, not structured blocks ──
_TEXT_FORMATS = frozenset({
    "json", "yaml", "yml", "txt", "md", "py", "js", "ts",
    "html", "xml", "css", "csv", "sh", "sql", "r", "lua",
    "toml", "ini", "cfg", "env", "log", "tex", "bib",
})

# ── Structured formats: content as list-of-dicts is intentional ──
_STRUCTURED_FORMATS = frozenset({"docx", "pdf", "pptx", "xlsx"})


def _is_model_wrapper_block(block: dict) -> bool:
    """
    Detect OpenRouter/ChatGPT structured wrapper blocks.
    These wrap actual content inside {"type": "raw"|"text"|"code", "text": "..."}.
    
    Distinguishes from legitimate DOCX/PDF/PPTX blocks which use types like
    "paragraph", "heading", "table", "list", "image", "title", "slide", etc.
    """
    if not isinstance(block, dict) or "text" not in block:
        return False
    block_type = block.get("type", "")
    # These are the wrapper types models use, NOT legitimate content block types
    return block_type in ("raw", "text", "code", "content", "")


def _unwrap_content_blocks(content_val, format_type: str):
    """
    Normalize content from various model output formats into what the
    format-specific creators expect.
    
    Handles:
      1. Stringified JSON: '["paragraph text"]' → parsed list
      2. OpenRouter wrapper blocks: [{"type":"raw","text":"..."}] → raw string
      3. Single-element list of string: ["just text"] → "just text"
      4. Pass-through for already-correct formats
    """
    # ── Step 1: Parse stringified JSON ──
    # Models sometimes send content as a JSON string instead of parsed object
    if isinstance(content_val, str):
        stripped = content_val.strip()
        if (stripped.startswith("[") and stripped.endswith("]")) or \
           (stripped.startswith("{") and stripped.endswith("}")):
            try:
                parsed = json.loads(stripped)
                # Only accept if it parsed to list/dict (not a JSON number/string)
                if isinstance(parsed, (list, dict)):
                    log.info(f"Content normalization: parsed stringified JSON "
                             f"({type(parsed).__name__}, format={format_type})")
                    content_val = parsed
            except (json.JSONDecodeError, ValueError):
                pass  # Not valid JSON, keep as string

    # ── Step 2: Handle list content ──
    if isinstance(content_val, list) and content_val:

        # Check if ALL items are model wrapper blocks
        if all(_is_model_wrapper_block(b) for b in content_val if isinstance(b, dict)):
            all_dicts = all(isinstance(b, dict) for b in content_val)

            if all_dicts and format_type in _TEXT_FORMATS:
                # Text format: unwrap and join text fields
                texts = [b.get("text", "") for b in content_val]
                unwrapped = "\n".join(t for t in texts if t)
                log.info(f"B5 unwrap: extracted text from {len(content_val)} "
                         f"wrapper blocks for format={format_type}")
                return unwrapped

            elif all_dicts and format_type in _STRUCTURED_FORMATS:
                # Structured format but with wrapper blocks instead of real content blocks
                # Try to extract and re-parse the text content
                texts = [b.get("text", "") for b in content_val]
                combined = "\n".join(t for t in texts if t)
                
                # For DOCX: if the extracted text looks like markdown, let _create_word handle it
                if format_type == "docx":
                    log.info(f"B5 unwrap: converting {len(content_val)} wrapper blocks "
                             f"to string for DOCX markdown→structured conversion")
                    return combined
                
                # For others: try to parse as JSON (model might have embedded structured data)
                try:
                    parsed = json.loads(combined)
                    if isinstance(parsed, list):
                        log.info(f"B5 unwrap: parsed embedded JSON from wrapper blocks "
                                 f"for format={format_type}")
                        return parsed
                except (json.JSONDecodeError, ValueError):
                    pass
                
                return combined

        # Single-element string list: ["just the content"]
        if len(content_val) == 1 and isinstance(content_val[0], str):
            if format_type in _TEXT_FORMATS:
                log.info(f"Content normalization: unwrapped single-element string list")
                return content_val[0]

    # ── Step 3: Dict content for text formats ──
    # Model sent content={"key": "value"} for a JSON file instead of a string
    if isinstance(content_val, dict) and format_type in _TEXT_FORMATS:
        if format_type in ("json", "yaml", "yml"):
            # Will be serialized later in the else branch, pass through
            pass
        else:
            # Other text formats: stringify
            log.info(f"Content normalization: converting dict to string for format={format_type}")
            return json.dumps(content_val, indent=2, ensure_ascii=False)

    return content_val


@mcp.tool()
async def create_file(
    # Accept EITHER nested data OR flat params
    data: dict = None,
    format: str = None,
    filename: str = None,
    content: Any = None,
    title: str = None,
    slides_data: list = None,
    persistent: bool = PERSISTENT_FILES
) -> dict:
    """
    Create a file. Accepts BOTH formats:
    - Nested: {"data": {"format": "html", "filename": "test.html", "content": "..."}}
    - Flat: {"format": "html", "filename": "test.html", "content": "..."}
    
    Formats: pdf, docx, pptx, xlsx, csv, txt, html, xml, py, js, json, md, etc.
    """
    log.debug("Creating file via tool")
    
    # ✅ AUTO-WRAPPER: Detect and normalize input format
    if data is not None:
        # Already wrapped - use as-is
        actual_data = data
    else:
        # Flat params provided - auto-wrap them
        actual_data = {}
        if format is not None:
            actual_data['format'] = format
        if filename is not None:
            actual_data['filename'] = filename
        if content is not None:
            actual_data['content'] = content
        if title is not None:
            actual_data['title'] = title
        if slides_data is not None:
            actual_data['slides_data'] = slides_data
    
    # Now use actual_data for all operations
    folder_path = _generate_unique_folder()
    format_type = (actual_data.get("format") or "").lower()
    filename_val = actual_data.get("filename")
    content_val = actual_data.get("content")
    title_val = actual_data.get("title")

    # ── B5 + Content Normalization ──
    # Unwrap structured wrapper blocks, parse stringified JSON, etc.
    # Must run BEFORE format routing so each creator gets clean input.
    content_val = _unwrap_content_blocks(content_val, format_type)

    if format_type == "pdf":
        result = _create_pdf(
            content_val if isinstance(content_val, list) else [str(content_val or "")], 
            filename_val, 
            folder_path=folder_path
        )
    elif format_type == "pptx":
        slides = actual_data.get("slides_data") or []
        if not slides and content_val:
            # Models often send PPTX content as structured blocks via `content`
            # rather than the `slides_data` parameter. Convert automatically.
            if isinstance(content_val, list):
                slides = _content_blocks_to_slides(content_val, title=title_val)
            elif isinstance(content_val, str):
                slides = [{"title": title_val or "", "content": [{"type": "paragraph", "text": content_val}]}]
        result = _create_presentation(
            slides, 
            filename_val, 
            folder_path=folder_path, 
            title=title_val
        )
    elif format_type == "docx":
        result = _create_word(
            content_val if content_val is not None else [], 
            filename_val, 
            folder_path=folder_path, 
            title=title_val
        )
    elif format_type == "xlsx":
        normalized_content = _normalize_xlsx_content(content_val) if content_val is not None else []
        result = _create_excel(
            normalized_content, 
            filename_val, 
            folder_path=folder_path, 
            title=title_val
        )
    elif format_type == "csv":
        result = _create_csv(
            content_val if content_val is not None else [], 
            filename_val, 
            folder_path=folder_path
        )
    else:
        use_filename = filename_val or f"export.{format_type or 'txt'}"
        # If content is a dict/list and format is JSON/YAML, serialize properly
        raw_content = content_val
        if isinstance(raw_content, (dict, list)):
            if format_type in ("json", "yaml", "yml"):
                raw_content = json.dumps(raw_content, indent=2, ensure_ascii=False)
            else:
                raw_content = str(raw_content)
        result = _create_raw_file(
            raw_content if raw_content is not None else "", 
            use_filename, 
            folder_path=folder_path
        )

    if not persistent:
        _cleanup_files(folder_path, FILES_DELAY)

    return {"url": result["url"]}

@mcp.tool()
async def generate_and_archive(files_data: list[dict], archive_format: str = "zip", archive_name: str = None, persistent: bool = PERSISTENT_FILES) -> dict:
    """files_data=[{"format":"pdf","filename":"report.pdf","content":[{"type":"title","text":"..."},{"type":"paragraph","text":"..."}],"title":"..."},{"format":"docx","filename":"doc.docx","content":[{"type":"title","text":"..."},{"type":"list","items":[...]}],"title":"..."},{"format":"pptx","filename":"slides.pptx","slides_data":[{"title":"...","content":[...],"image_query":"...","image_position":"left|right|top|bottom","image_size":"small|medium|large"}],"title":"..."},{"format":"xlsx","filename":"data.xlsx","content":[["Header1","Header2"],["Val1","Val2"]],"title":"..."},{"format":"csv","filename":"data.csv","content":[[...]]},{"format":"txt|xml|py|etc","filename":"file.ext","content":"string"}]"""
    log.debug("Generating archive via tool")
    folder_path = _generate_unique_folder()
    generated_paths: list[str] = []

    for file_info in files_data or []:
        fmt = (file_info.get("format") or "").lower()
        fname = file_info.get("filename")
        content = file_info.get("content")
        title = file_info.get("title")

        try:
            if fmt == "pdf":
                res = _create_pdf(content if isinstance(content, list) else [str(content or "")], fname, folder_path=folder_path)
            elif fmt == "pptx":
                res = _create_presentation(file_info.get("slides_data", []), fname, folder_path=folder_path, title=title)
            elif fmt == "docx":
                res = _create_word(content if content is not None else [], fname, folder_path=folder_path, title=title)
            elif fmt == "xlsx":
                normalized = _normalize_xlsx_content(content) if content is not None else []
                res = _create_excel(normalized, fname, folder_path=folder_path, title=title)
            elif fmt == "csv":
                res = _create_csv(content if content is not None else [], fname, folder_path=folder_path)
            else:
                use_fname = fname or f"export.{fmt or 'txt'}"
                res = _create_raw_file(content if content is not None else "", use_fname, folder_path=folder_path)
        except Exception as e:
            log.error(f"Error generating file {fname or '<no name>'}: {e}", exc_info=True)
            raise

        generated_paths.append(res["path"])

    timestamp = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
    archive_basename = f"{archive_name or 'archive'}_{timestamp}"
    archive_filename = f"{archive_basename}.zip" if archive_format.lower() not in ("7z", "tar.gz") else f"{archive_basename}.{archive_format}"
    archive_path = os.path.join(folder_path, archive_filename)

    if archive_format.lower() == "7z":
        with py7zr.SevenZipFile(archive_path, mode='w') as archive:
            for p in generated_paths:
                archive.write(p, os.path.relpath(p, folder_path))
    elif archive_format.lower() == "tar.gz":
        with tarfile.open(archive_path, "w:gz") as tar:
            for p in generated_paths:
                tar.add(p, arcname=os.path.relpath(p, folder_path))
    else:
        with zipfile.ZipFile(archive_path, 'w') as zipf:
            for p in generated_paths:
                zipf.write(p, os.path.relpath(p, folder_path))

    if not persistent:
        _cleanup_files(folder_path, FILES_DELAY)

    return {"url": _public_url(folder_path, archive_filename)}

from sse_starlette.sse import EventSourceResponse

class SimpleRequestContext:
    def __init__(self, request):
        self.request = request

class SimpleCtx:
    def __init__(self, request):
        self.request_context = SimpleRequestContext(request)

async def handle_sse(request: Request) -> Response:
    """Handle SSE transport for MCP - supports both GET and POST"""
    
    if request.method == "POST":
        try:
            message = await request.json()
            log.debug(f"Received POST message: {message}")
            
            response = {
                "jsonrpc": "2.0",
                "id": message.get("id"),
                "result": None
            }
            
            method = message.get("method")
            
            if method == "initialize":
                response["result"] = {
                    "protocolVersion": "2024-11-05",
                    "capabilities": {
                        "tools": {},
                        "logging": {}
                    },
                    "serverInfo": {
                        "name": "file_export_mcp",
                        "version": SCRIPT_VERSION
                    }
                }
            elif method == "tools/list":
                response["result"] = {
                    "tools": [
                        {
                            "name": "create_file",
                            "description": "Create files in various formats (pdf, docx, pptx, xlsx, csv, txt, xml, py, etc.). Supports rich content including titles, paragraphs, lists, tables, images via queries, and more.",
                            "inputSchema": {
                                "type": "object",
                                "properties": {
                                    "data": {
                                        "type": "object",
                                        "description": "File data configuration",
                                        "properties": {
                                            "format": {
                                                "type": "string",
                                                "enum": ["pdf", "docx", "pptx", "xlsx", "csv", "txt", "xml", "py", "json", "md"],
                                                "description": "Output file format"
                                            },
                                            "filename": {
                                                "type": "string",
                                                "description": "Name of the file to create (optional, will be auto-generated if not provided)"
                                            },
                                            "title": {
                                                "type": "string",
                                                "description": "Document title (for docx, pptx, xlsx, pdf)"
                                            },
                                            "content": {
                                                "description": "Content varies by format. For pdf/docx: array (objects or strings). For xlsx/csv: 2D array. For pptx: use slides_data instead. For txt/xml/py: string",
                                                "oneOf": [
                                                    {
                                                        "type": "array",
                                                        "items": {
                                                            "anyOf": [
                                                                { "type": "string" },
                                                                { "type": "number" },
                                                                { "type": "boolean" },
                                                                { "type": "object" },
                                                                {
                                                                    "type": "array",
                                                                    "items": {
                                                                        "anyOf": [
                                                                            { "type": "string" },
                                                                            { "type": "number" },
                                                                            { "type": "boolean" },
                                                                            { "type": "object" },
                                                                            { "type": "null" }
                                                                        ]
                                                                    }
                                                                }
                                                            ]
                                                        }
                                                    },
                                                    { "type": "string" },
                                                    { "type": "object" },
                                                    { "type": "null" }
                                                ]
                                            },
                                            "slides_data": {
                                                "type": "array",
                                                "description": "For pptx format only: array of slide objects",
                                                "items": {
                                                    "type": "object",
                                                    "properties": {
                                                        "title": {"type": "string"},
                                                        "content": {
                                                            "type": "array",
                                                            "items": {"type": "string"}
                                                        },
                                                        "image_query": {
                                                            "type": "string",
                                                            "description": "Search query for image (Unsplash, Pexels, or local SD)"
                                                        },
                                                        "image_position": {
                                                            "type": "string",
                                                            "enum": ["left", "right", "top", "bottom"],
                                                            "description": "Position of the image on the slide"
                                                        },
                                                        "image_size": {
                                                            "type": "string",
                                                            "enum": ["small", "medium", "large"],
                                                            "description": "Size of the image"
                                                        }
                                                    }
                                                }
                                            }
                                        },
                                        "required": ["format"]
                                    },
                                    "persistent": {
                                        "type": "boolean",
                                        "description": "Whether to keep files permanently (default: false, files deleted after delay)"
                                    }
                                },
                                "required": ["data"]
                            }
                        },
                        {
                            "name": "generate_and_archive",
                            "description": "Generate multiple files and create an archive (zip, 7z, tar.gz)",
                            "inputSchema": {
                                "type": "object",
                                "properties": {
                                    "files_data": {
                                        "type": "array",
                                        "description": "Array of file data objects",
                                        "items": {
                                            "type": "object",
                                            "properties": {
                                                "format": { "type": "string" },
                                                "filename": { "type": "string" },
                                                "content": {
                                                    "description": "For pdf/docx: array (objects or strings). For xlsx/csv: 2D array. For others: string/object",
                                                    "oneOf": [
                                                        {
                                                            "type": "array",
                                                            "items": {
                                                                "anyOf": [
                                                                    { "type": "string" },
                                                                    { "type": "number" },
                                                                    { "type": "boolean" },
                                                                    { "type": "object" },
                                                                    {
                                                                        "type": "array",
                                                                        "items": {
                                                                            "anyOf": [
                                                                                { "type": "string" },
                                                                                { "type": "number" },
                                                                                { "type": "boolean" },
                                                                                { "type": "object" },
                                                                                { "type": "null" }
                                                                            ]
                                                                        }
                                                                    }
                                                                ]
                                                            }
                                                        },
                                                        { "type": "string" },
                                                        { "type": "object" },
                                                        { "type": "null" }
                                                    ]
                                                },
                                                "title": { "type": "string" },
                                                "slides_data": {
                                                    "type": "array",
                                                    "description": "For pptx format only: array of slide objects",
                                                    "items": {
                                                        "type": "object",
                                                        "properties": {
                                                            "title": { "type": "string" },
                                                            "content": {
                                                                "type": "array",
                                                                "items": { "type": "string" }
                                                            },
                                                            "image_query": {
                                                                "type": "string",
                                                                "description": "Search query for image (Unsplash, Pexels, or local SD)"
                                                            },
                                                            "image_position": {
                                                                "type": "string",
                                                                "enum": ["left", "right", "top", "bottom"],
                                                                "description": "Position of the image on the slide"
                                                            },
                                                            "image_size": {
                                                                "type": "string",
                                                                "enum": ["small", "medium", "large"],
                                                                "description": "Size of the image"
                                                            }
                                                        }
                                                    }
                                                }
                                            },
                                            "required": ["format"]
                                        }
                                    },
                                    "archive_format": {"type": "string", "enum": ["zip", "7z", "tar.gz"]},
                                    "archive_name": {"type": "string"},
                                    "persistent": {"type": "boolean"}
                                },
                                "required": ["files_data"]
                            }
                        },
                        {
                            "name": "full_context_document",
                            "description": "Extract and return the complete structure, content, and metadata of a document (docx, xlsx, pptx). Returns a JSON structure with indexed elements (paragraphs, headings, tables, cells, slides, images) that can be referenced for editing or review.",
                            "inputSchema": {
                                "type": "object",
                                "properties": {
                                    "file_id": {
                                        "type": "string",
                                        "description": "The file ID from OpenWebUI file upload"
                                    },
                                    "file_name": {
                                        "type": "string",
                                        "description": "The name of the file with extension (e.g., 'report.docx', 'data.xlsx', 'presentation.pptx')"
                                    }
                                },
                                "required": ["file_id", "file_name"]
                            }
                        },
                        {
                            "name": "edit_document",
                            "description": "Edit an existing document (docx, xlsx, pptx) using structured operations. Supports inserting/deleting elements and updating content. ALWAYS call full_context_document() first to get proper IDs and references. Preserves formatting and returns a download link for the edited file.",
                            "inputSchema": {
                                "type": "object",
                                "properties": {
                                    "file_id": {
                                        "type": "string",
                                        "description": "The file ID from OpenWebUI"
                                    },
                                    "file_name": {
                                        "type": "string",
                                        "description": "The name of the file with extension"
                                    },
                                    "edits": {
                                        "type": "object",
                                        "description": "Edit operations and content changes",
                                        "properties": {
                                            "ops": {
                                                "type": "array",
                                                "description": "Structural operations (insert/delete). For PPTX: ['insert_after', slide_id, 'nK'], ['insert_before', slide_id, 'nK'], ['delete_slide', slide_id]. For DOCX: ['insert_after', para_xml_id, 'nK'], ['insert_before', para_xml_id, 'nK'], ['delete_paragraph', para_xml_id]. For XLSX: ['insert_row', 'sheet_name', row_idx], ['delete_row', 'sheet_name', row_idx], ['insert_column', 'sheet_name', col_idx], ['delete_column', 'sheet_name', col_idx]",
                                                "items": {
                                                    "type": "array",
                                                    "items": {
                                                        "oneOf": [
                                                            {"type": "string"},
                                                            {"type": "integer"}
                                                        ]
                                                    }
                                                }
                                            },
                                            "content_edits": {
                                                "type": "array",
                                                "description": "Content updates. Prefer object items: {target, value}. For PPTX: target 'sid:<slide_id>/shid:<shape_id>' or 'nK:slot:title'/'body'/'table'. For DOCX: 'pid:<para_xml_id>' or 'tid:<table_xml_id>/cid:<cell_xml_id>' or 'nK'. For XLSX: 'A1', 'B5'.",
                                                "items": {
                                                    "type": "object",
                                                    "required": ["target", "value"],
                                                    "properties": {
                                                        "target": {
                                                            "type": "string",
                                                            "description": "Target reference (element ID or cell ref)"
                                                        },
                                                        "value": {
                                                            "description": "New content (string, number, boolean, array of strings, or 2D array for tables)",
                                                            "oneOf": [
                                                                {"type": "string"},
                                                                {"type": "number"},
                                                                {"type": "boolean"},
                                                                {"type": "array", "items": {"type": "string"}},
                                                                {
                                                                    "type": "array",
                                                                    "items": {
                                                                        "type": "array",
                                                                        "items": {
                                                                            "oneOf": [
                                                                                {"type": "string"},
                                                                                {"type": "number"},
                                                                                {"type": "boolean"},
                                                                                {"type": "null"}
                                                                            ]
                                                                        }
                                                                    }
                                                                }
                                                            ]
                                                        }
                                                    },
                                                    "additionalProperties": False
                                                }
                                            }
                                        }
                                    }
                                },
                                "required": ["file_id", "file_name", "edits"]
                            }
                        },
                        {
                            "name": "review_document",
                            "description": "Review and add comments/corrections to an existing document (docx, xlsx, pptx). Returns a download link for the reviewed document with comments added. For Excel, the index MUST be a cell reference (e.g., 'A1', 'B5', 'C10') as returned by full_context_document. For Word: use either an integer paragraph index or 'pid:<para_xml_id>'. For PowerPoint: use either an integer slide index or 'sid:<slide_id>' (optionally 'sid:<slide_id>/shid:<shape_id>' to target a shape).",
                            "inputSchema": {
                                "type": "object",
                                "properties": {
                                    "file_id": {
                                        "type": "string",
                                        "description": "The file ID from OpenWebUI"
                                    },
                                    "file_name": {
                                        "type": "string",
                                        "description": "The name of the file with extension"
                                    },
                                            "review_comments": {
                                                "type": "array",
                                                "description": "Array of objects {index, comment}. For Excel: index must be a cell reference string like 'A1', 'B3'. For Word: integer paragraph index or 'pid:<para_xml_id>'. For PowerPoint: integer slide index or 'sid:<slide_id>' (optionally 'sid:<slide_id>/shid:<shape_id>').",
                                                "items": {
                                                    "type": "object",
                                                    "required": ["index", "comment"],
                                                    "properties": {
                                                        "index": {
                                                            "description": "Index/reference: For Excel use cell reference (e.g., 'A1'); for Word/PowerPoint use integer or an id key string like 'pid:<para_xml_id>' / 'sid:<slide_id>'",
                                                            "oneOf": [
                                                                {"type": "string"},
                                                                {"type": "integer"}
                                                            ]
                                                        },
                                                        "comment": {
                                                            "type": "string",
                                                            "description": "Comment or correction text"
                                                        }
                                                    },
                                                    "additionalProperties": False
                                                }
                                            }
                                },
                                "required": ["file_id", "file_name", "review_comments"]
                            }
                        }
                    ]
                }
            elif method == "tools/call":
                params = message.get("params", {})
                tool_name = params.get("name")
                arguments = params.get("arguments", {}) or {}
                ctx = SimpleCtx(request)

                try:
                    if tool_name == "create_file":
                        result = await create_file(**arguments)
                        response["result"] = {
                            "content": [
                                {
                                    "type": "text",
                                    "text": f"File created successfully: {result.get('url')}"
                                }
                            ],
                            "isError": False
                        }

                    elif tool_name == "generate_and_archive":
                        result = await generate_and_archive(**arguments)
                        response["result"] = {
                            "content": [
                                {
                                    "type": "text",
                                    "text": f"Archive created successfully: {result.get('url')}"
                                }
                            ],
                            "isError": False
                        }

                    elif tool_name == "full_context_document":
                        arguments.setdefault("ctx", ctx)
                        result = await full_context_document(**arguments)
                        response["result"] = {
                            "content": [
                                {
                                    "type": "text",
                                    "text": result
                                }
                            ],
                            "isError": False
                        }

                    elif tool_name == "edit_document":
                        arguments.setdefault("ctx", ctx)
                        # Normalize 'edits.content_edits' to list of [target, value] for backward compatibility
                        try:
                            edits_arg = arguments.get("edits")
                            if isinstance(edits_arg, dict):
                                ce = edits_arg.get("content_edits")
                                if isinstance(ce, list) and (len(ce) == 0 or isinstance(ce[0], dict)):
                                    edits_arg["content_edits"] = [
                                        [item.get("target"), item.get("value")]
                                        for item in (ce or [])
                                        if isinstance(item, dict) and "target" in item and "value" in item
                                    ]
                        except Exception:
                            pass
                        result = await edit_document(**arguments)
                        response["result"] = {
                            "content": [
                                {
                                    "type": "text",
                                    "text": json.dumps(result, indent=2, ensure_ascii=False)
                                }
                            ],
                            "isError": False
                        }

                    elif tool_name == "review_document":
                        arguments.setdefault("ctx", ctx)
                        # Normalize 'review_comments' to list of [index, comment] for backward compatibility
                        try:
                            rc = arguments.get("review_comments")
                            if isinstance(rc, list) and (len(rc) == 0 or isinstance(rc[0], dict)):
                                arguments["review_comments"] = [
                                    [item.get("index"), item.get("comment")]
                                    for item in (rc or [])
                                    if isinstance(item, dict) and "index" in item and "comment" in item
                                ]
                        except Exception:
                            pass
                        result = await review_document(**arguments)
                        response["result"] = {
                            "content": [
                                {
                                    "type": "text",
                                    "text": json.dumps(result, indent=2, ensure_ascii=False)
                                }
                            ],
                            "isError": False
                        }

                    else:
                        response["error"] = {
                            "code": -32601,
                            "message": f"Tool not found: {tool_name}"
                        }
                except Exception as e:
                    log.error(f"Error executing tool {tool_name}: {e}", exc_info=True)
                    response["result"] = {
                        "content": [
                            {
                                "type": "text",
                                "text": f"Error executing tool: {str(e)}"
                            }
                        ],
                        "isError": True
                    }
            else:
                response["error"] = {
                    "code": -32601,
                    "message": f"Method not found: {method}"
                }
            
            return JSONResponse(response)
            
        except Exception as e:
            log.error(f"Error handling POST request: {e}", exc_info=True)
            return JSONResponse(
                {
                    "jsonrpc": "2.0",
                    "id": None,
                    "error": {
                        "code": -32700,
                        "message": f"Parse error: {str(e)}"
                    }
                },
                status_code=400
            )
    
    else:
        async def event_generator():
            """Generator for SSE events with correct text format"""
            try:
                endpoint_data = json.dumps({"endpoint": "/sse"})
                yield f"event: endpoint\ndata: {endpoint_data}\n\n"
                
                import asyncio
                while True:
                    await asyncio.sleep(15)
                    yield f"event: ping\ndata: {{}}\n\n"
                    
            except asyncio.CancelledError:
                log.info("SSE connection closed by client")
                raise
            except Exception as e:
                log.error(f"SSE Error: {e}", exc_info=True)
                error_data = json.dumps({"error": str(e)})
                yield f"event: error\ndata: {error_data}\n\n"
        
        return EventSourceResponse(
            event_generator(),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "X-Accel-Buffering": "no",
                "Connection": "keep-alive"
            }
        )

async def handle_messages(request: Request) -> Response:
    """Handle POST requests to /messages endpoint"""
    try:
        data = await request.json()
        log.debug(f"Received message: {data}")
        
        response = {
            "jsonrpc": "2.0",
            "id": data.get("id"),
            "result": {
                "content": [
                    {
                        "type": "text",
                        "text": "Message received"
                    }
                ]
            }
        }
        
        return JSONResponse(response)
    except Exception as e:
        log.error(f"Message handling error: {e}", exc_info=True)
        return JSONResponse(
            {
                "jsonrpc": "2.0",
                "id": None,
                "error": {
                    "code": -32700,
                    "message": f"Parse error: {str(e)}"
                }
            },
            status_code=400
        )

async def health_check(request: Request) -> Response:
    """Health check endpoint"""
    return JSONResponse({"status": "healthy", "server": "file_export_mcp"})

app = Starlette(
    debug=True,
    routes=[
        Route("/sse", endpoint=handle_sse, methods=["GET", "POST"]),
        Route("/messages", endpoint=handle_messages, methods=["POST"]),
        Route("/health", endpoint=health_check, methods=["GET"]),
    ]
)

if __name__ == "__main__":
    # Default to "stdio" if MODE is not set, so it works with MetaMCP
    mode = os.getenv("MODE", "stdio").lower()

    if mode == "sse":
        port = int(os.getenv("MCP_HTTP_PORT", "9004"))
        host = os.getenv("MCP_HTTP_HOST", "0.0.0.0")
            
        log.info(f"Starting file_export_mcp version {SCRIPT_VERSION}")
        log.info(f"Starting file_export_mcp in SSE mode on http://{host}:{port}")
        
        uvicorn.run(
            app,
            host=host,
            port=port,
            access_log=False,
            log_level="info",
            use_colors=False
        )
    elif mode == "http":
        port = int(os.getenv("MCP_HTTP_PORT", "9004"))
        host = os.getenv("MCP_HTTP_HOST", "0.0.0.0")
        
        log.info(f"Starting file_export_mcp version {SCRIPT_VERSION}")
        log.info(f"Starting file_export_mcp in http mode on http://{host}:{port}")

        mcp.run(transport="streamable-http")
        
    else:
        # Fallback to STDIO (MetaMCP default)
        log.info(f"Starting file_export_mcp version {SCRIPT_VERSION} in STDIO mode")
        mcp.run(transport="stdio")